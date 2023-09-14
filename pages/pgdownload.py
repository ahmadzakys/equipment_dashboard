######-----Import Dash-----#####
import dash
from dash import dcc
from dash import html, callback
import dash_bootstrap_components as dbc
from dash.dependencies import Input, Output
from dash.exceptions import PreventUpdate

from pptx import Presentation
from pptx.util import Cm, Pt

import os

import pandas as pd
import numpy as np
import plotly.graph_objects as go
import plotly.io as pio
from plotly.subplots import make_subplots

from datetime import date
from dateutil.relativedelta import relativedelta

dash.register_page(__name__, name='Download')

##-----page

## -----LAYOUT-----
layout = html.Div([
                html.Br(),
                
                html.Div([
                    html.Button('Download Slides', 
                                      id='download_button', 
                                      n_clicks=0, 
                                      style={'fontSize': 15, 'color':'#2a3f5f','font-family':'Verdana','display': 'inline-block'}), 
                    dcc.Download(id='download'),
                    ]),

                html.Br(),
                html.P(id='download_date', style={'fontSize': 15, 'color':'#2a3f5f','font-family':'Verdana'}),

    html.Footer('ABL',
            style={'textAlign': 'center', 
                   'fontSize': 20, 
                   'background-color':'#2a3f5f',
                   'color':'white',
                   'position': 'fixed',
                    'bottom': '0',
                    'width': '100%'})

    ], style={
        'paddingLeft':'10px',
        'paddingRight':'10px',
    })

#### Callback Auto Update Chart & Data

@callback(
    [Output('download', 'data'),
     Output('download_date', 'children')],
    [Input('download_button', 'n_clicks'),
     Input('store', 'data')]
)
def update_charts(n_clicks, data):
    if n_clicks == 0:
        raise PreventUpdate
    ######################
    # Pre Processing
    ######################

    ###################### Sumatra ######################
    df_sumatra = pd.DataFrame(data['Bulk Sumatra'])
    df_sumatra = df_sumatra.replace(r'^\s*$', np.nan, regex=True)

    df_sumatra[['Last Overhaul GEN01', 'Last Overhaul GEN02', 'Last Overhaul GEN03', 'Last Overhaul GEN04']] = \
    df_sumatra[['Last Overhaul GEN01', 'Last Overhaul GEN02', 'Last Overhaul GEN03', 'Last Overhaul GEN04']].fillna(date.today())
    df_sumatra['Date'] = pd.to_datetime(df_sumatra['Date'], dayfirst = True)
    df_sumatra['Last Overhaul GEN01'] = pd.to_datetime(df_sumatra['Last Overhaul GEN01'], dayfirst = True)
    df_sumatra['Last Overhaul GEN02'] = pd.to_datetime(df_sumatra['Last Overhaul GEN02'], dayfirst = True)
    df_sumatra['Last Overhaul GEN03'] = pd.to_datetime(df_sumatra['Last Overhaul GEN03'], dayfirst = True)
    df_sumatra['Last Overhaul GEN04'] = pd.to_datetime(df_sumatra['Last Overhaul GEN04'], dayfirst = True)
    df_sumatra['Next Docking Intermediate Survey (IS)'] = pd.to_datetime(df_sumatra['Next Docking Intermediate Survey (IS)'], dayfirst = True)

    # Hydrolic Oil Utilization
    last = df_sumatra.iloc[-1]
    hou1 = (last['CRANE 1']-last['Last Rpl RH CR01'])/6000
    hou1 = round(hou1*100,0)
    hou1 = min(hou1, 100)

    hou2 = (last['CRANE 2']-last['Last Rpl RH CR02'])/6000
    hou2 = round(hou2*100,0)
    hou2 = min(hou2, 100)

    # Wire Rope
    wire_rope_c1 = [['HW RHOL', last['HW RHOL CR01'], last['CRANE 1']],
                    ['HW LHOL', last['HW LHOL CR01'], last['CRANE 1']],
                    ['CW RHOL', last['CW RHOL CR01'], last['CRANE 1']],
                    ['CW LHOL', last['CW LHOL CR01'], last['CRANE 1']]]
    wire_rope_c1 = pd.DataFrame(wire_rope_c1, columns=['Crane 1', 'Last Replacement', 'Current RH'])
    wire_rope_c1['Utilization RH'] = wire_rope_c1['Current RH']-wire_rope_c1['Last Replacement']
    wire_rope_c1['Utilization %'] = round(wire_rope_c1['Utilization RH']/4000*100,0)
    util1 = []
    for i in wire_rope_c1['Utilization %']:
        util1.append(min(i, 100))
    wire_rope_c1['Utilization'] = util1
    wire_rope_c1['Max'] = 100

    wire_rope_c2 = [['HW RHOL', last['HW RHOL CR02'], last['CRANE 2']],
                    ['HW LHOL', last['HW LHOL CR02'], last['CRANE 2']],
                    ['CW RHOL', last['CW RHOL CR02'], last['CRANE 2']],
                    ['CW LHOL', last['CW LHOL CR02'], last['CRANE 2']]]
    wire_rope_c2 = pd.DataFrame(wire_rope_c2, columns=['Crane 2', 'Last Replacement', 'Current RH'])
    wire_rope_c2['Utilization RH'] = wire_rope_c2['Current RH']-wire_rope_c2['Last Replacement']
    wire_rope_c2['Utilization %'] = round(wire_rope_c2['Utilization RH']/4000*100,0)
    util2 = []
    for i in wire_rope_c2['Utilization %']:
        util2.append(min(i, 100))
    wire_rope_c2['Utilization'] = util2
    wire_rope_c2['Max'] = 100

    # Engine GOH
    engine_goh = [['DG1', last['DG1'], last['Last TOH DG1'], last['Last GOH DG1']],
                ['DG2', last['DG2'], last['Last TOH DG2'], last['Last GOH DG2']],
                ['DG3', last['DG3'], last['Last TOH DG3'], last['Last GOH DG3']],
                ['EDG', last['EDG'], last['Last TOH EDG'], last['Last GOH EDG']]]
    engine_goh = pd.DataFrame(engine_goh, columns=['Sumatra', 'Current HR', 'Last TOH', 'Last GOH'])
    engine_goh['Last TOH'] = engine_goh['Last TOH'].fillna(0)
    engine_goh['Utilization TOH'] = engine_goh['Current HR']-engine_goh['Last TOH']
    engine_goh['Utilization GOH'] = engine_goh['Current HR']-engine_goh['Last GOH']
    engine_goh['Hours till TOH'] = 12000-engine_goh['Utilization TOH']
    engine_goh['Hours till GOH'] = 24000-engine_goh['Utilization GOH']
    engine_goh['TOH %'] = round((1-engine_goh['Hours till TOH']/12000)*100,0)
    engine_goh['GOH %'] = round((1-engine_goh['Hours till GOH']/24000)*100,0)

    TOH = []
    for i in engine_goh['TOH %']:
        TOH.append(min(i, 100))
    engine_goh['TOH'] = TOH

    GOH = []
    for i in engine_goh['GOH %']:
        GOH.append(min(i, 100))
    engine_goh['GOH'] = GOH

    # Conveyor Belt Replacement
    conveyor_br = [['BC1', last['BC1'], last['Last RPL HR BC1']],
                ['BC2', last['BC2'], last['Last RPL HR BC2']],
                ['BC3', last['BC3'], last['Last RPL HR BC3']],
                ['BC4', last['BC4'], last['Last RPL HR BC4']],
                ['BC5', last['BC5'], last['Last RPL HR BC5']],
                ['SHL1', last['SHL1'], last['Last RPL HR SHL1']],
                ['SHL2', last['SHL2'], last['Last RPL HR SHL2']]]
    conveyor_br = pd.DataFrame(conveyor_br, columns=['Sumatra', 'Current RH', 'Last Rpl HR'])
    conveyor_br['Utilization (hrs)'] = conveyor_br['Current RH']-conveyor_br['Last Rpl HR']
    conveyor_br['Utilization %'] = [round(conveyor_br['Utilization (hrs)'][0]/8000*100,0),
                                round(conveyor_br['Utilization (hrs)'][1]/8000*100,0),
                                round(conveyor_br['Utilization (hrs)'][2]/5000*100,0),
                                round(conveyor_br['Utilization (hrs)'][3]/5000*100,0),
                                round(conveyor_br['Utilization (hrs)'][4]/5000*100,0),
                                round(conveyor_br['Utilization (hrs)'][5]/5000*100,0),
                                round(conveyor_br['Utilization (hrs)'][6]/5000*100,0),]

    util = []
    for i in conveyor_br['Utilization %']:
        util.append(min(i, 100))
    conveyor_br['Utilization'] = util
    conveyor_br['Max'] = 100

    # Alternator OH
    alternator_oh = [['GEN01', last['Last Overhaul GEN01'],],
                    ['GEN02', last['Last Overhaul GEN02'],],
                    ['GEN03', last['Last Overhaul GEN03'],],
                    ['GEN04', last['Last Overhaul GEN04'],]]
    alternator_oh = pd.DataFrame(alternator_oh, columns=['Sumatra', 'Last Overhaul (Date)',])

    last_oh = []
    for i in alternator_oh['Last Overhaul (Date)']:
        last_oh.append(i+relativedelta(months=60))
    alternator_oh['Next Overhaul (Date)'] = last_oh

    alternator_oh['Todays Date'] = date.today()
    alternator_oh['Todays Date'] = alternator_oh['Todays Date'].astype('datetime64[ns]')
    alternator_oh['Time to Go'] = (alternator_oh['Next Overhaul (Date)'] - alternator_oh['Todays Date']).dt.days.astype('int64')
    alternator_oh['Five Years in Days'] = 1827
    alternator_oh['Utilization %'] = round((alternator_oh['Five Years in Days']-alternator_oh['Time to Go'])/ \
                                    alternator_oh['Five Years in Days']*100,0)

    util = []
    for i in alternator_oh['Utilization %']:
        util.append(min(i, 100))
    alternator_oh['Utilization'] = util
    alternator_oh['Max'] = 100

    #Critical Spare
    crane = last['Crane']
    df_crane = pd.DataFrame({'names' : ['progress',' '],
                    'values' :  [crane, 100 - crane]})

    conveyor = last['Conveyor']
    df_conveyor = pd.DataFrame({'names' : ['progress',' '],
                    'values' :  [conveyor, 100 - conveyor]})

    engine = last['Engine']
    df_engine = pd.DataFrame({'names' : ['progress',' '],
                    'values' :  [engine, 100 - engine]})


    def shape_cu(value):
        shape=''
        if value == 'Green': shape="img/green check.png"
        elif value == 'Yellow': shape="img/yellow cross.png"
        elif value == 'Red': shape="img/red cross.png"
        elif value == 'Gray': shape="img/gray blank.png"
        return(shape)

    def shape_cs(value):
        shape=''
        if value<=(50) : shape="img/red cross.png"
        elif value<=(80) : shape="img/yellow cross.png"
        elif value<=100 : shape="img/green check.png"
        else : shape="img/gray blank.png"
        return(shape)

    def col_day_delta(value):
        col=''
        if value<=180 : col="img/red calendar.png"
        elif value<=365 : col="img/yellow calendar.png"
        else: col="img/green calendar.png"
        return(col)

    path = os.getcwd()
    download_date = html.Strong('Downloaded slides as per ' + str(last['Date'].strftime("%d %b %Y")))

    t_sumatra = last['Next Docking Intermediate Survey (IS)']
    ConMakerShape_sumatra = shape_cu(last['Conveyor-Maker Check-Up'])
    CraneMakerShape_sumatra = shape_cu(last['Crane-Maker Check-Up'])
    EngMakerShape_sumatra = shape_cu(last['Engine-Maker Check-Up'])
    OthersShape_sumatra = shape_cu(last['Others'])
    NextDockingCol_sumatra = col_day_delta((last['Next Docking Intermediate Survey (IS)'] - pd.Timestamp.today()).days)
    ConSpareShape_sumatra = shape_cu(last['Conveyor Belt Spares'])
    CraneSpareShape_sumatra = shape_cu(last['Crane-Wire Rope Spares']) 
    CraneShape_sumatra = shape_cs(last['Crane'])
    ConShape_sumatra = shape_cs(last['Conveyor'])
    EngShape_sumatra = shape_cs(last['Engine'])


    #Chart
    sumatra_hou_c1 = plot_hou1(hou1)
    sumatra_hou_c2= plot_hou2(hou2)
    sumatra_toh = plot_toh(engine_goh['TOH'])
    sumatra_goh = plot_goh(engine_goh['GOH'])
    sumatra_crane_1 = plot_wr1(wire_rope_c1['Utilization'], wire_rope_c1)
    sumatra_crane_2 = plot_wr2(wire_rope_c2['Utilization'], wire_rope_c2)
    sumatra_cbu = plot_cbr(conveyor_br['Utilization'], 'Sumatra', conveyor_br)
    sumatra_au = plot_aoh(alternator_oh['Utilization'], 'Sumatra', alternator_oh)
    sumatra_cs_crane = plot_cs(df_crane)
    sumatra_cs_conveyor = plot_cs(df_conveyor)
    sumatra_cs_engine = plot_cs(df_engine)

    #Save Img 
    pio.write_image(sumatra_hou_c1, 'img/hou_c1 Bulk Sumatra.png')
    pio.write_image(sumatra_hou_c2, 'img/hou_c2 Bulk Sumatra.png')
    pio.write_image(sumatra_toh, 'img/toh Bulk Sumatra.png')
    pio.write_image(sumatra_goh, 'img/goh Bulk Sumatra.png')
    pio.write_image(sumatra_crane_1, 'img/crane_1 Bulk Sumatra.png')
    pio.write_image(sumatra_crane_2, 'img/crane_2 Bulk Sumatra.png')
    pio.write_image(sumatra_cbu, 'img/cbu Bulk Sumatra.png')
    pio.write_image(sumatra_au, 'img/au Bulk Sumatra.png')
    pio.write_image(sumatra_cs_crane, 'img/cs_crane Bulk Sumatra.png')
    pio.write_image(sumatra_cs_conveyor, 'img/cs_conveyor Bulk Sumatra.png')
    pio.write_image(sumatra_cs_engine, 'img/cs_engine Bulk Sumatra.png')

    ###################### Derawan ######################
    df_derawan = pd.DataFrame(data['Bulk Derawan'])
    df_derawan = df_derawan.replace(r'^\s*$', np.nan, regex=True)

    df_derawan[['Last Overhaul GEN01', 'Last Overhaul GEN02', 'Last Overhaul GEN03', 'Last Overhaul GEN04']] = \
    df_derawan[['Last Overhaul GEN01', 'Last Overhaul GEN02', 'Last Overhaul GEN03', 'Last Overhaul GEN04']].fillna(date.today())
    df_derawan['Last Overhaul GEN01'] = pd.to_datetime(df_derawan['Last Overhaul GEN01'], dayfirst = True)
    df_derawan['Last Overhaul GEN02'] = pd.to_datetime(df_derawan['Last Overhaul GEN02'], dayfirst = True)
    df_derawan['Last Overhaul GEN03'] = pd.to_datetime(df_derawan['Last Overhaul GEN03'], dayfirst = True)
    df_derawan['Last Overhaul GEN04'] = pd.to_datetime(df_derawan['Last Overhaul GEN04'], dayfirst = True)
    df_derawan['Next Docking Intermediate Survey (IS)'] = pd.to_datetime(df_derawan['Next Docking Intermediate Survey (IS)'], dayfirst = True)

    df_derawan['Last RPL date Feeder'] = pd.to_datetime(df_derawan['Last RPL date Feeder'], dayfirst = True)
    df_derawan['Last RPL date SHL'] = pd.to_datetime(df_derawan['Last RPL date SHL'], dayfirst = True)
    last = df_derawan.iloc[-1]

    # Wire Rope
    wire_rope_c1 = [['HW RHOL', last['HW RHOL CR01'], last['CRANE 1']],
                    ['HW LHOL', last['HW LHOL CR01'], last['CRANE 1']],
                    ['CW RHOL', last['CW RHOL CR01'], last['CRANE 1']],
                    ['CW LHOL', last['CW LHOL CR01'], last['CRANE 1']],
                    ['EW RHOL', last['EW RHOL CR01'], last['CRANE 1']]]
    wire_rope_c1 = pd.DataFrame(wire_rope_c1, columns=['Crane 1', 'Last Replacement', 'Current RH'])
    wire_rope_c1['Utilization RH'] = wire_rope_c1['Current RH']-wire_rope_c1['Last Replacement']
    wire_rope_c1['Utilization %'] = [round(wire_rope_c1['Utilization RH'][0]/2500*100,0),
                                    round(wire_rope_c1['Utilization RH'][1]/2500*100,0),
                                    round(wire_rope_c1['Utilization RH'][2]/2500*100,0),
                                    round(wire_rope_c1['Utilization RH'][3]/2500*100,0),
                                    round(wire_rope_c1['Utilization RH'][4]/1200*100,0),]
    util1 = []
    for i in wire_rope_c1['Utilization %']:
        util1.append(min(i, 100))
    wire_rope_c1['Utilization'] = util1
    wire_rope_c1['Max'] = 100


    wire_rope_c2 = [['HW RHOL', last['HW RHOL CR02'], last['CRANE 2']],
                    ['HW LHOL', last['HW LHOL CR02'], last['CRANE 2']],
                    ['CW RHOL', last['CW RHOL CR02'], last['CRANE 2']],
                    ['CW LHOL', last['CW LHOL CR02'], last['CRANE 2']],
                    ['EW RHOL', last['EW RHOL CR02'], last['CRANE 2']]]
    wire_rope_c2 = pd.DataFrame(wire_rope_c2, columns=['Crane 2', 'Last Replacement', 'Current RH'])
    wire_rope_c2['Utilization RH'] = wire_rope_c2['Current RH']-wire_rope_c2['Last Replacement']
    wire_rope_c2['Utilization %'] = [round(wire_rope_c2['Utilization RH'][0]/2500*100,0),
                                    round(wire_rope_c2['Utilization RH'][1]/2500*100,0),
                                    round(wire_rope_c2['Utilization RH'][2]/2500*100,0),
                                    round(wire_rope_c2['Utilization RH'][3]/2500*100,0),
                                    round(wire_rope_c2['Utilization RH'][4]/1200*100,0),]
    util2 = []
    for i in wire_rope_c2['Utilization %']:
        util2.append(min(i, 100))
    wire_rope_c2['Utilization'] = util2
    wire_rope_c2['Max'] = 100


    # Engine GOH
    engine_goh = [['DG1', last['DG1'], last['Last TOH DG1'], last['Last GOH DG1']],
                ['DG2', last['DG2'], last['Last TOH DG2'], last['Last GOH DG2']],
                ['DG3', last['DG3'], last['Last TOH DG3'], last['Last GOH DG3']],
                ['EDG', last['AE'], last['Last TOH EDG'], last['Last GOH EDG']]]
    engine_goh = pd.DataFrame(engine_goh, columns=['Derawan', 'Current HR', 'Last TOH', 'Last GOH'])
    engine_goh['Last TOH'] = engine_goh['Last TOH'].fillna(0)
    engine_goh['Utilization TOH'] = engine_goh['Current HR']-engine_goh['Last TOH']
    engine_goh['Utilization GOH'] = engine_goh['Current HR']-engine_goh['Last GOH']
    engine_goh['Hours till TOH'] = 12000-engine_goh['Utilization TOH']
    engine_goh['Hours till GOH'] = 24000-engine_goh['Utilization GOH']
    engine_goh['TOH %'] = round((1-engine_goh['Hours till TOH']/12000)*100,0)
    engine_goh['GOH %'] = round((1-engine_goh['Hours till GOH']/24000)*100,0)

    TOH = []
    for i in engine_goh['TOH %']:
        TOH.append(min(i, 100))
    engine_goh['TOH'] = TOH

    GOH = []
    for i in engine_goh['GOH %']:
        GOH.append(min(i, 100))
    engine_goh['GOH'] = GOH

    # Conveyor Belt Replacement
    conveyor_br = [['Feeder', last['Last RPL date Feeder']],
                ['SHL', last['Last RPL date SHL']]]
    conveyor_br = pd.DataFrame(conveyor_br, columns=['Derawan', 'Last RPL date'])
    conveyor_br['Last RPL date'] = conveyor_br['Last RPL date'].astype('datetime64[ns]')
    conveyor_br['Todays Date'] = date.today()
    conveyor_br['Todays Date'] = conveyor_br['Todays Date'].astype('datetime64[ns]')
    conveyor_br['Time to Go (days)'] = (conveyor_br['Todays Date'] - conveyor_br['Last RPL date']).dt.days.astype('int64')
    conveyor_br['Time to Go (months)'] = round(conveyor_br['Time to Go (days)']/30,2)
    conveyor_br['Utilization %'] = round(conveyor_br['Time to Go (months)']/35*100,0)

    util = []
    for i in conveyor_br['Utilization %']:
        util.append(min(i, 100))
    conveyor_br['Utilization'] = util
    conveyor_br['Max'] = 100

    # Alternator OH
    alternator_oh = [['GEN01', last['Last Overhaul GEN01'],],
                    ['GEN02', last['Last Overhaul GEN02'],],
                    ['GEN03', last['Last Overhaul GEN03'],],
                    ['GEN04', last['Last Overhaul GEN04'],]]
    alternator_oh = pd.DataFrame(alternator_oh, columns=['Derawan', 'Last Overhaul (Date)',])

    last_oh = []
    for i in alternator_oh['Last Overhaul (Date)']:
        last_oh.append(i+relativedelta(months=60))
    alternator_oh['Next Overhaul (Date)'] = last_oh

    alternator_oh['Todays Date'] = date.today()
    alternator_oh['Todays Date'] = alternator_oh['Todays Date'].astype('datetime64[ns]')
    alternator_oh['Time to Go'] = (alternator_oh['Next Overhaul (Date)'] - alternator_oh['Todays Date']).dt.days.astype('int64')
    alternator_oh['Five Years in Days'] = 1827
    alternator_oh['Utilization %'] = round((alternator_oh['Five Years in Days']-alternator_oh['Time to Go'])/ \
                                    alternator_oh['Five Years in Days']*100,0)

    util = []
    for i in alternator_oh['Utilization %']:
        util.append(min(i, 100))
    alternator_oh['Utilization'] = util
    alternator_oh['Max'] = 100

    #Critical Spare
    conveyor = last['Conveyor']
    df_conveyor = pd.DataFrame({'names' : ['progress',' '],
                    'values' :  [conveyor, 100 - conveyor]})
    
    def shape_cu(value):
        shape=''
        if value == 'Green': shape="img/green check.png"
        elif value == 'Yellow': shape="img/yellow cross.png"
        elif value == 'Red': shape="img/red cross.png"
        elif value == 'Gray': shape="img/gray blank.png"
        return(shape)

    def shape_cs(value):
        shape=''
        if value<=(50) : shape="img/red cross.png"
        elif value<=(80) : shape="img/yellow cross.png"
        elif value<=100 : shape="img/green check.png"
        else : shape="img/gray blank.png"
        return(shape)

    def col_day_delta(value):
        col=''
        if value<=180 : col="img/red calendar.png"
        elif value<=365 : col="img/yellow calendar.png"
        else: col="img/green calendar.png"
        return(col)

    t_derawan = last['Next Docking Intermediate Survey (IS)']
    ConMakerShape_derawan = shape_cu(last['Conveyor-Maker Check-Up'])
    CraneMakerShape_derawan = shape_cu(last['Crane-Maker Check-Up'])
    EngMakerShape_derawan = shape_cu(last['Engine-Maker Check-Up'])
    OthersShape_derawan = shape_cu(last['Others'])
    NextDockingCol_derawan = col_day_delta((last['Next Docking Intermediate Survey (IS)'] - pd.Timestamp.today()).days)
    ConSpareShape_derawan = shape_cu(last['Conveyor Belt Spares'])
    CraneSpareShape_derawan = shape_cu(last['Crane-Wire Rope Spares']) 
    CraneShape_derawan = shape_cs(last['Crane'])
    ConShape_derawan = shape_cs(last['Conveyor'])
    EngShape_derawan = shape_cs(last['Engine'])

    #Chart
    derawan_toh = plot_toh(engine_goh['TOH'])
    derawan_goh = plot_goh(engine_goh['GOH'])
    derawan_crane_1 = plot_wr1_dk(wire_rope_c1['Utilization'], wire_rope_c1)
    derawan_crane_2 = plot_wr2_dk(wire_rope_c2['Utilization'], wire_rope_c2)
    derawan_cbu = plot_cbr(conveyor_br['Utilization'], 'Derawan', conveyor_br)
    derawan_au = plot_aoh(alternator_oh['Utilization'], 'Derawan', alternator_oh)
    derawan_cs_conveyor = plot_cs(df_conveyor)

    #Save Img
    pio.write_image(derawan_toh, 'img/toh Bulk Derawan.png')
    pio.write_image(derawan_goh, 'img/goh Bulk Derawan.png')
    pio.write_image(derawan_crane_1, 'img/crane_1 Bulk Derawan.png')
    pio.write_image(derawan_crane_2, 'img/crane_2 Bulk Derawan.png')
    pio.write_image(derawan_cbu, 'img/cbu Bulk Derawan.png')
    pio.write_image(derawan_au, 'img/au Bulk Derawan.png')
    pio.write_image(derawan_cs_conveyor, 'img/cs_conveyor Bulk Derawan.png')


    ###################### Karimun ######################
    df_karimun = pd.DataFrame(data['Bulk Karimun'])
    df_karimun = df_karimun.replace(r'^\s*$', np.nan, regex=True)

    df_karimun[['Last Overhaul GEN01', 'Last Overhaul GEN02', 'Last Overhaul GEN03', 'Last Overhaul GEN04']] = \
    df_karimun[['Last Overhaul GEN01', 'Last Overhaul GEN02', 'Last Overhaul GEN03', 'Last Overhaul GEN04']].fillna(date.today())
    df_karimun['Last Overhaul GEN01'] = pd.to_datetime(df_karimun['Last Overhaul GEN01'], dayfirst = True)
    df_karimun['Last Overhaul GEN02'] = pd.to_datetime(df_karimun['Last Overhaul GEN02'], dayfirst = True)
    df_karimun['Last Overhaul GEN03'] = pd.to_datetime(df_karimun['Last Overhaul GEN03'], dayfirst = True)
    df_karimun['Last Overhaul GEN04'] = pd.to_datetime(df_karimun['Last Overhaul GEN04'], dayfirst = True)
    df_karimun['Next Docking Intermediate Survey (IS)'] = pd.to_datetime(df_karimun['Next Docking Intermediate Survey (IS)'], dayfirst = True)

    last = df_karimun.iloc[-1]

    # Wire Rope
    wire_rope_c1 = [['HW RHOL', last['HW RHOL CR01'], last['CRANE 1']],
                    ['HW LHOL', last['HW LHOL CR01'], last['CRANE 1']],
                    ['GW RHOL', last['GW RHOL CR01'], last['CRANE 1']],
                    ['GW LHOL', last['GW LHOL CR01'], last['CRANE 1']],
                    ['EW RHOL', last['EW RHOL CR01'], last['CRANE 1']]]
    wire_rope_c1 = pd.DataFrame(wire_rope_c1, columns=['Crane 1', 'Last Replacement', 'Current RH'])
    wire_rope_c1['Utilization RH'] = wire_rope_c1['Current RH']-wire_rope_c1['Last Replacement']
    wire_rope_c1['Utilization %'] = [round(wire_rope_c1['Utilization RH'][0]/2000*100,0),
                                    round(wire_rope_c1['Utilization RH'][1]/2000*100,0),
                                    round(wire_rope_c1['Utilization RH'][2]/2000*100,0),
                                    round(wire_rope_c1['Utilization RH'][3]/2000*100,0),
                                    round(wire_rope_c1['Utilization RH'][4]/1500*100,0),]
    util1 = []
    for i in wire_rope_c1['Utilization %']:
        util1.append(min(i, 100))
    wire_rope_c1['Utilization'] = util1
    wire_rope_c1['Max'] = 100


    wire_rope_c2 = [['HW RHOL', last['HW RHOL CR02'], last['CRANE 2']],
                    ['HW LHOL', last['HW LHOL CR02'], last['CRANE 2']],
                    ['GW RHOL', last['GW RHOL CR02'], last['CRANE 2']],
                    ['GW LHOL', last['GW LHOL CR02'], last['CRANE 2']],
                    ['EW RHOL', last['EW RHOL CR02'], last['CRANE 2']]]
    wire_rope_c2 = pd.DataFrame(wire_rope_c2, columns=['Crane 2', 'Last Replacement', 'Current RH'])
    wire_rope_c2['Utilization RH'] = wire_rope_c2['Current RH']-wire_rope_c2['Last Replacement']
    wire_rope_c2['Utilization %'] = [round(wire_rope_c2['Utilization RH'][0]/2000*100,0),
                                    round(wire_rope_c2['Utilization RH'][1]/2000*100,0),
                                    round(wire_rope_c2['Utilization RH'][2]/2000*100,0),
                                    round(wire_rope_c2['Utilization RH'][3]/2000*100,0),
                                    round(wire_rope_c2['Utilization RH'][4]/1500*100,0),]
    util2 = []
    for i in wire_rope_c2['Utilization %']:
        util2.append(min(i, 100))
    wire_rope_c2['Utilization'] = util2
    wire_rope_c2['Max'] = 100


    # Engine GOH
    engine_goh = [['DG1', last['MG1'], last['Last TOH DG1'], last['Last GOH DG1']],
                ['DG2', last['MG2'], last['Last TOH DG2'], last['Last GOH DG2']],
                ['DG3', last['MG3'], last['Last TOH DG3'], last['Last GOH DG3']],
                ['EDG', last['HG'], last['Last TOH EDG'], last['Last GOH EDG']]]
    engine_goh = pd.DataFrame(engine_goh, columns=['Karimun', 'Current HR', 'Last TOH', 'Last GOH'])
    engine_goh['Last TOH'] = engine_goh['Last TOH'].fillna(0)
    # engine_goh['Utilization TOH'] = 12000-engine_goh['Utilization GOH']
    # engine_goh['Utilization GOH'] = 24000-engine_goh['Utilization GOH']
    engine_goh['Hours till TOH'] = 12000-engine_goh['Last GOH']
    engine_goh['Hours till GOH'] = 24000-engine_goh['Last GOH']
    engine_goh['TOH %'] = round((1-engine_goh['Hours till TOH']/12000)*100,0)
    engine_goh['GOH %'] = round((1-engine_goh['Hours till GOH']/24000)*100,0)

    TOH = []
    for i in engine_goh['TOH %']:
        TOH.append(min(i, 100))
    engine_goh['TOH'] = TOH

    GOH = []
    for i in engine_goh['GOH %']:
        GOH.append(min(i, 100))
    engine_goh['GOH'] = GOH

    # Conveyor Belt Replacement
    conveyor_br = [['FB01', last['FEEDER 1'], last['Last RPL HR FB01']],
                ['FB02', last['FEEDER 2'], last['Last RPL HR FB02']],
                ['CB1', last['CB1'], last['Last RPL HR CB1']],
                ['CB2', last['CB2'], last['Last RPL HR CB2']],
                ['CB3', last['CB3'], last['Last RPL HR CB3']],
                ['CB4', last['CB4'], last['Last RPL HR CB4']],
                ['SHL1', last['SHL1'], last['Last RPL HR SHL1']]]
    conveyor_br = pd.DataFrame(conveyor_br, columns=['Karimun', 'Current RH', 'Last Rpl HR'])
    conveyor_br['Utilization (hrs)'] = conveyor_br['Current RH']-conveyor_br['Last Rpl HR']
    conveyor_br['Utilization %'] = [round(conveyor_br['Utilization (hrs)'][0]/5000*100,0),
                                round(conveyor_br['Utilization (hrs)'][1]/5000*100,0),
                                round(conveyor_br['Utilization (hrs)'][2]/10000*100,0),
                                round(conveyor_br['Utilization (hrs)'][3]/10000*100,0),
                                round(conveyor_br['Utilization (hrs)'][4]/10000*100,0),
                                round(conveyor_br['Utilization (hrs)'][5]/5000*100,0),
                                round(conveyor_br['Utilization (hrs)'][6]/10000*100,0),]

    util = []
    for i in conveyor_br['Utilization %']:
        util.append(min(i, 100))
    conveyor_br['Utilization'] = util
    conveyor_br['Max'] = 100

    # Alternator OH
    alternator_oh = [['GEN01', last['Last Overhaul GEN01'],],
                    ['GEN02', last['Last Overhaul GEN02'],],
                    ['GEN03', last['Last Overhaul GEN03'],],
                    ['GEN04', last['Last Overhaul GEN04'],]]
    alternator_oh = pd.DataFrame(alternator_oh, columns=['Karimun', 'Last Overhaul (Date)',])

    last_oh = []
    for i in alternator_oh['Last Overhaul (Date)']:
        last_oh.append(i+relativedelta(months=60))
    alternator_oh['Next Overhaul (Date)'] = last_oh

    alternator_oh['Todays Date'] = date.today()
    alternator_oh['Todays Date'] = alternator_oh['Todays Date'].astype('datetime64[ns]')
    alternator_oh['Time to Go'] = (alternator_oh['Next Overhaul (Date)'] - alternator_oh['Todays Date']).dt.days.astype('int64')
    alternator_oh['Five Years in Days'] = 1827
    alternator_oh['Utilization %'] = round((alternator_oh['Five Years in Days']-alternator_oh['Time to Go'])/ \
                                    alternator_oh['Five Years in Days']*100,0)

    util = []
    for i in alternator_oh['Utilization %']:
        util.append(min(i, 100))
    alternator_oh['Utilization'] = util
    alternator_oh['Max'] = 100

    #Critical Spare
    conveyor = last['Conveyor']
    df_conveyor = pd.DataFrame({'names' : ['progress',' '],
                    'values' :  [conveyor, 100 - conveyor]})
    
    def shape_cu(value):
        shape=''
        if value == 'Green': shape="img/green check.png"
        elif value == 'Yellow': shape="img/yellow cross.png"
        elif value == 'Red': shape="img/red cross.png"
        elif value == 'Gray': shape="img/gray blank.png"
        return(shape)

    def shape_cs(value):
        shape=''
        if value<=(50) : shape="img/red cross.png"
        elif value<=(80) : shape="img/yellow cross.png"
        elif value<=100 : shape="img/green check.png"
        else : shape="img/gray blank.png"
        return(shape)

    def col_day_delta(value):
        col=''
        if value<=180 : col="img/red calendar.png"
        elif value<=365 : col="img/yellow calendar.png"
        else: col="img/green calendar.png"
        return(col)

    t_karimun = last['Next Docking Intermediate Survey (IS)']
    ConMakerShape_karimun = shape_cu(last['Conveyor-Maker Check-Up'])
    CraneMakerShape_karimun = shape_cu(last['Crane-Maker Check-Up'])
    EngMakerShape_karimun = shape_cu(last['Engine-Maker Check-Up'])
    OthersShape_karimun = shape_cu(last['Others'])
    NextDockingCol_karimun = col_day_delta((last['Next Docking Intermediate Survey (IS)'] - pd.Timestamp.today()).days)
    ConSpareShape_karimun = shape_cu(last['Conveyor Belt Spares'])
    CraneSpareShape_karimun = shape_cu(last['Crane-Wire Rope Spares']) 
    CraneShape_karimun = shape_cs(last['Crane'])
    ConShape_karimun = shape_cs(last['Conveyor'])
    EngShape_karimun = shape_cs(last['Engine'])

    #Chart
    karimun_toh = plot_toh(engine_goh['TOH'])
    karimun_goh = plot_goh(engine_goh['GOH'])
    karimun_crane_1 = plot_wr1_dk(wire_rope_c1['Utilization'], wire_rope_c1)
    karimun_crane_2 = plot_wr2_dk(wire_rope_c2['Utilization'], wire_rope_c2)
    karimun_cbu = plot_cbr(conveyor_br['Utilization'], 'Karimun', conveyor_br)
    karimun_au = plot_aoh(alternator_oh['Utilization'], 'Karimun', alternator_oh)
    karimun_cs_conveyor = plot_cs(df_conveyor)

    #Save Img
    pio.write_image(karimun_toh, 'img/toh Bulk Karimun.png')
    pio.write_image(karimun_goh, 'img/goh Bulk Karimun.png')
    pio.write_image(karimun_crane_1, 'img/crane_1 Bulk Karimun.png')
    pio.write_image(karimun_crane_2, 'img/crane_2 Bulk Karimun.png')
    pio.write_image(karimun_cbu, 'img/cbu Bulk Karimun.png')
    pio.write_image(karimun_au, 'img/au Bulk Karimun.png')
    pio.write_image(karimun_cs_conveyor, 'img/cs_conveyor Bulk Karimun.png')

    ###################### Dewata ######################
    df_dewata = pd.DataFrame(data['Bulk Dewata'])
    df_dewata = df_dewata.replace(r'^\s*$', np.nan, regex=True)

    df_dewata[['Last Overhaul GEN01', 'Last Overhaul GEN02', 'Last Overhaul GEN03', 'Last Overhaul GEN04']] = \
    df_dewata[['Last Overhaul GEN01', 'Last Overhaul GEN02', 'Last Overhaul GEN03', 'Last Overhaul GEN04']].fillna(date.today())
    df_dewata['Last Overhaul GEN01'] = pd.to_datetime(df_dewata['Last Overhaul GEN01'], dayfirst = True)
    df_dewata['Last Overhaul GEN02'] = pd.to_datetime(df_dewata['Last Overhaul GEN02'], dayfirst = True)
    df_dewata['Last Overhaul GEN03'] = pd.to_datetime(df_dewata['Last Overhaul GEN03'], dayfirst = True)
    df_dewata['Last Overhaul GEN04'] = pd.to_datetime(df_dewata['Last Overhaul GEN04'], dayfirst = True)
    df_dewata['Next Docking Intermediate Survey (IS)'] = pd.to_datetime(df_dewata['Next Docking Intermediate Survey (IS)'], dayfirst = True)


    # Hydrolic Oil Utilization
    last = df_dewata.iloc[-1]
    hou1 = (last['CRANE 1 (PMS 500)']-last['Last Rpl RH CR01'])/6000
    hou1 = round(hou1*100,0)
    hou1 = min(hou1, 100)

    hou2 = (last['CRANE 2 (PMS 500)']-last['Last Rpl RH CR02'])/6000
    hou2 = round(hou2*100,0)
    hou2 = min(hou2, 100)

    # Wire Rope
    wire_rope_c1 = [['HW RHOL', last['HW RHOL CR01'], last['CRANE 1 (PMS 500)']],
                    ['HW LHOL', last['HW LHOL CR01'], last['CRANE 1 (PMS 500)']],
                    ['CW RHOL', last['CW RHOL CR01'], last['CRANE 1 (PMS 500)']],
                    ['CW LHOL', last['CW LHOL CR01'], last['CRANE 1 (PMS 500)']]]
    wire_rope_c1 = pd.DataFrame(wire_rope_c1, columns=['Crane 1', 'Last Replacement', 'Current RH'])
    wire_rope_c1['Utilization RH'] = wire_rope_c1['Current RH']-wire_rope_c1['Last Replacement']
    wire_rope_c1['Utilization %'] = round(wire_rope_c1['Utilization RH']/4000*100,0)
    util1 = []
    for i in wire_rope_c1['Utilization %']:
        util1.append(min(i, 100))
    wire_rope_c1['Utilization'] = util1
    wire_rope_c1['Max'] = 100

    wire_rope_c2 = [['HW RHOL', last['HW RHOL CR02'], last['CRANE 2 (PMS 500)']],
                    ['HW LHOL', last['HW LHOL CR02'], last['CRANE 2 (PMS 500)']],
                    ['CW RHOL', last['CW RHOL CR02'], last['CRANE 2 (PMS 500)']],
                    ['CW LHOL', last['CW LHOL CR02'], last['CRANE 2 (PMS 500)']]]
    wire_rope_c2 = pd.DataFrame(wire_rope_c2, columns=['Crane 2', 'Last Replacement', 'Current RH'])
    wire_rope_c2['Utilization RH'] = wire_rope_c2['Current RH']-wire_rope_c2['Last Replacement']
    wire_rope_c2['Utilization %'] = round(wire_rope_c2['Utilization RH']/4000*100,0)
    util2 = []
    for i in wire_rope_c2['Utilization %']:
        util2.append(min(i, 100))
    wire_rope_c2['Utilization'] = util2
    wire_rope_c2['Max'] = 100

    # Engine GOH
    engine_goh = [['DG1', last['DG1'], last['Last TOH DG1'], last['Last GOH DG1']],
                ['DG2', last['DG2'], last['Last TOH DG2'], last['Last GOH DG2']],
                ['DG3', last['DG3'], last['Last TOH DG3'], last['Last GOH DG3']],
                ['EDG', last['DG4'], last['Last TOH EDG'], last['Last GOH EDG']]]
    engine_goh = pd.DataFrame(engine_goh, columns=['Dewata', 'Current HR', 'Last TOH', 'Last GOH'])
    engine_goh['Last TOH'] = engine_goh['Last TOH'].fillna(0)
    # engine_goh['Utilization TOH'] = engine_goh['Current HR']-engine_goh['Last TOH']
    # engine_goh['Utilization GOH'] = engine_goh['Current HR']-engine_goh['Last GOH']
    engine_goh['Hours till TOH'] = 12000-engine_goh['Last GOH']
    engine_goh['Hours till GOH'] = 24000-engine_goh['Last GOH']
    engine_goh['TOH %'] = round((1-engine_goh['Hours till TOH']/12000)*100,0)
    engine_goh['GOH %'] = round((1-engine_goh['Hours till GOH']/24000)*100,0)

    TOH = []
    for i in engine_goh['TOH %']:
        TOH.append(min(i, 100))
    engine_goh['TOH'] = TOH

    GOH = []
    for i in engine_goh['GOH %']:
        GOH.append(min(i, 100))
    engine_goh['GOH'] = GOH

    # Conveyor Belt Replacement
    conveyor_br = [['FB01', last['FB1'], last['Last RPL HR FB01']],
                ['FB02', last['FB2'], last['Last RPL HR FB02']],
                ['BC01', last['BC1'], last['Last RPL HR BC1']],
                ['BC02', last['BC2'], last['Last RPL HR BC2']],
                ['BC03', last['BC3'], last['Last RPL HR BC3']],
                ['SHL', last['SHL1'], last['Last RPL HR SHL']],]
    conveyor_br = pd.DataFrame(conveyor_br, columns=['Dewata', 'Current RH', 'Last Rpl HR'])
    conveyor_br['Utilization (hrs)'] = conveyor_br['Current RH']-conveyor_br['Last Rpl HR']
    conveyor_br['Utilization %'] = [round(conveyor_br['Utilization (hrs)'][0]/10000*100,0),
                                round(conveyor_br['Utilization (hrs)'][1]/10000*100,0),
                                round(conveyor_br['Utilization (hrs)'][2]/10000*100,0),
                                round(conveyor_br['Utilization (hrs)'][3]/10000*100,0),
                                round(conveyor_br['Utilization (hrs)'][4]/10000*100,0),
                                round(conveyor_br['Utilization (hrs)'][5]/10000*100,0),]

    util = []
    for i in conveyor_br['Utilization %']:
        util.append(min(i, 100))
    conveyor_br['Utilization'] = util
    conveyor_br['Max'] = 100

    # Alternator OH
    alternator_oh = [['GEN01', last['Last Overhaul GEN01'],],
                    ['GEN02', last['Last Overhaul GEN02'],],
                    ['GEN03', last['Last Overhaul GEN03'],],
                    ['GEN04', last['Last Overhaul GEN04'],]]
    alternator_oh = pd.DataFrame(alternator_oh, columns=['Dewata', 'Last Overhaul (Date)',])

    last_oh = []
    for i in alternator_oh['Last Overhaul (Date)']:
        last_oh.append(i+relativedelta(months=60))
    alternator_oh['Next Overhaul (Date)'] = last_oh

    alternator_oh['Todays Date'] = date.today()
    alternator_oh['Todays Date'] = alternator_oh['Todays Date'].astype('datetime64[ns]')
    alternator_oh['Time to Go'] = (alternator_oh['Next Overhaul (Date)'] - alternator_oh['Todays Date']).dt.days.astype('int64')
    alternator_oh['Five Years in Days'] = 1827
    alternator_oh['Utilization %'] = round((alternator_oh['Five Years in Days']-alternator_oh['Time to Go'])/ \
                                    alternator_oh['Five Years in Days']*100,0)

    util = []
    for i in alternator_oh['Utilization %']:
        util.append(min(i, 100))
    alternator_oh['Utilization'] = util
    alternator_oh['Max'] = 100

    #Critical Spare
    crane = last['Crane']
    df_crane = pd.DataFrame({'names' : ['progress',' '],
                    'values' :  [crane, 100 - crane]})

    conveyor = last['Conveyor']
    df_conveyor = pd.DataFrame({'names' : ['progress',' '],
                    'values' :  [conveyor, 100 - conveyor]})

    engine = last['Engine']
    df_engine = pd.DataFrame({'names' : ['progress',' '],
                    'values' :  [engine, 100 - engine]})
    
    def shape_cu(value):
        shape=''
        if value == 'Green': shape="img/green check.png"
        elif value == 'Yellow': shape="img/yellow cross.png"
        elif value == 'Red': shape="img/red cross.png"
        elif value == 'Gray': shape="img/gray blank.png"
        return(shape)

    def shape_cs(value):
        shape=''
        if value<=(50) : shape="img/red cross.png"
        elif value<=(80) : shape="img/yellow cross.png"
        elif value<=100 : shape="img/green check.png"
        else : shape="img/gray blank.png"
        return(shape)

    def col_day_delta(value):
        col=''
        if value<=180 : col="img/red calendar.png"
        elif value<=365 : col="img/yellow calendar.png"
        else: col="img/green calendar.png"
        return(col)

    t_dewata = last['Next Docking Intermediate Survey (IS)']
    ConMakerShape_dewata = shape_cu(last['Conveyor-Maker Check-Up'])
    CraneMakerShape_dewata = shape_cu(last['Crane-Maker Check-Up'])
    EngMakerShape_dewata = shape_cu(last['Engine-Maker Check-Up'])
    OthersShape_dewata = shape_cu(last['Others'])
    NextDockingCol_dewata = col_day_delta((last['Next Docking Intermediate Survey (IS)'] - pd.Timestamp.today()).days)
    ConSpareShape_dewata = shape_cu(last['Conveyor Belt Spares'])
    CraneSpareShape_dewata = shape_cu(last['Crane-Wire Rope Spares']) 
    CraneShape_dewata = shape_cs(last['Crane'])
    ConShape_dewata = shape_cs(last['Conveyor'])
    EngShape_dewata = shape_cs(last['Engine'])

    #Chart
    dewata_hou_c1 = plot_hou1(hou1)
    dewata_hou_c2= plot_hou2(hou2)
    dewata_toh = plot_toh(engine_goh['TOH'])
    dewata_goh = plot_goh(engine_goh['GOH'])
    dewata_crane_1 = plot_wr1(wire_rope_c1['Utilization'], wire_rope_c1)
    dewata_crane_2 = plot_wr2(wire_rope_c2['Utilization'], wire_rope_c2)
    dewata_cbu = plot_cbr(conveyor_br['Utilization'], 'Dewata', conveyor_br)
    dewata_au = plot_aoh(alternator_oh['Utilization'], 'Dewata', alternator_oh)
    dewata_cs_crane = plot_cs(df_crane)
    dewata_cs_conveyor = plot_cs(df_conveyor)
    dewata_cs_engine = plot_cs(df_engine)

    #Save Img
    pio.write_image(dewata_hou_c1, 'img/hou_c1 Bulk Dewata.png')
    pio.write_image(dewata_hou_c2, 'img/hou_c2 Bulk Dewata.png')
    pio.write_image(dewata_toh, 'img/toh Bulk Dewata.png')
    pio.write_image(dewata_goh, 'img/goh Bulk Dewata.png')
    pio.write_image(dewata_crane_1, 'img/crane_1 Bulk Dewata.png')
    pio.write_image(dewata_crane_2, 'img/crane_2 Bulk Dewata.png')
    pio.write_image(dewata_cbu, 'img/cbu Bulk Dewata.png')
    pio.write_image(dewata_au, 'img/au Bulk Dewata.png')
    pio.write_image(dewata_cs_crane, 'img/cs_crane Bulk Dewata.png')
    pio.write_image(dewata_cs_conveyor, 'img/cs_conveyor Bulk Dewata.png')
    pio.write_image(dewata_cs_engine, 'img/cs_engine Bulk Dewata.png')

    ###################### Sumba ######################
    df_sumba = pd.DataFrame(data['Bulk Sumba'])
    df_sumba = df_sumba.replace(r'^\s*$', np.nan, regex=True)

    df_sumba[['Last Overhaul GEN01', 'Last Overhaul GEN02', 'Last Overhaul GEN03', 'Last Overhaul GEN04']] = \
    df_sumba[['Last Overhaul GEN01', 'Last Overhaul GEN02', 'Last Overhaul GEN03', 'Last Overhaul GEN04']].fillna(date.today())
    df_sumba['Last Overhaul GEN01'] = pd.to_datetime(df_sumba['Last Overhaul GEN01'], dayfirst = True)
    df_sumba['Last Overhaul GEN02'] = pd.to_datetime(df_sumba['Last Overhaul GEN02'], dayfirst = True)
    df_sumba['Last Overhaul GEN03'] = pd.to_datetime(df_sumba['Last Overhaul GEN03'], dayfirst = True)
    df_sumba['Last Overhaul GEN04'] = pd.to_datetime(df_sumba['Last Overhaul GEN04'], dayfirst = True)
    df_sumba['Next Docking Intermediate Survey (IS)'] = pd.to_datetime(df_sumba['Next Docking Intermediate Survey (IS)'], dayfirst = True)


    # Hydrolic Oil Utilization
    last = df_sumba.iloc[-1]
    hou1 = (last['CRANE 1 (DAILY TANK HYDROLICK)']-last['Last Rpl RH CR01'])/6000
    hou1 = round(hou1*100,0)
    hou1 = min(hou1, 100)

    hou2 = (last['CRANE 2 (DAILY TANK HYDROLICK)']-last['Last Rpl RH CR02'])/6000
    hou2 = round(hou2*100,0)
    hou2 = min(hou2, 100)

    # Wire Rope
    wire_rope_c1 = [['HW RHOL', last['HW RHOL CR01'], last['CRANE 1 (DAILY TANK HYDROLICK)']],
                    ['HW LHOL', last['HW LHOL CR01'], last['CRANE 1 (DAILY TANK HYDROLICK)']],
                    ['CW RHOL', last['CW RHOL CR01'], last['CRANE 1 (DAILY TANK HYDROLICK)']],
                    ['CW LHOL', last['CW LHOL CR01'], last['CRANE 1 (DAILY TANK HYDROLICK)']]]
    wire_rope_c1 = pd.DataFrame(wire_rope_c1, columns=['Crane 1', 'Last Replacement', 'Current RH'])
    wire_rope_c1['Utilization RH'] = wire_rope_c1['Current RH']-wire_rope_c1['Last Replacement']
    wire_rope_c1['Utilization %'] = round(wire_rope_c1['Utilization RH']/4000*100,0)
    util1 = []
    for i in wire_rope_c1['Utilization %']:
        util1.append(min(i, 100))
    wire_rope_c1['Utilization'] = util1
    wire_rope_c1['Max'] = 100

    wire_rope_c2 = [['HW RHOL', last['HW RHOL CR02'], last['CRANE 2 (DAILY TANK HYDROLICK)']],
                    ['HW LHOL', last['HW LHOL CR02'], last['CRANE 2 (DAILY TANK HYDROLICK)']],
                    ['CW RHOL', last['CW RHOL CR02'], last['CRANE 2 (DAILY TANK HYDROLICK)']],
                    ['CW LHOL', last['CW LHOL CR02'], last['CRANE 2 (DAILY TANK HYDROLICK)']]]
    wire_rope_c2 = pd.DataFrame(wire_rope_c2, columns=['Crane 2', 'Last Replacement', 'Current RH'])
    wire_rope_c2['Utilization RH'] = wire_rope_c2['Current RH']-wire_rope_c2['Last Replacement']
    wire_rope_c2['Utilization %'] = round(wire_rope_c2['Utilization RH']/4000*100,0)
    util2 = []
    for i in wire_rope_c2['Utilization %']:
        util2.append(min(i, 100))
    wire_rope_c2['Utilization'] = util2
    wire_rope_c2['Max'] = 100

    # Engine GOH
    engine_goh = [['DG1', last['DG1'], last['Last TOH DG1'], last['Last GOH DG1']],
                ['DG2', last['DG2'], last['Last TOH DG2'], last['Last GOH DG2']],
                ['DG3', last['DG3'], last['Last TOH DG3'], last['Last GOH DG3']],
                ['EDG', last['EDG'], last['Last TOH EDG'], last['Last GOH EDG']]]
    engine_goh = pd.DataFrame(engine_goh, columns=['Sumba', 'Current HR', 'Last TOH', 'Last GOH'])
    engine_goh['Last TOH'] = engine_goh['Last TOH'].fillna(0)
    # engine_goh['Utilization TOH'] = engine_goh['Current HR']-engine_goh['Last TOH']
    # engine_goh['Utilization GOH'] = engine_goh['Current HR']-engine_goh['Last GOH']
    engine_goh['Hours till TOH'] = 12000-engine_goh['Last GOH']
    engine_goh['Hours till GOH'] = 24000-engine_goh['Last GOH']
    engine_goh['TOH %'] = round((1-engine_goh['Hours till TOH']/12000)*100,0)
    engine_goh['GOH %'] = round((1-engine_goh['Hours till GOH']/24000)*100,0)

    TOH = []
    for i in engine_goh['TOH %']:
        TOH.append(min(i, 100))
    engine_goh['TOH'] = TOH

    GOH = []
    for i in engine_goh['GOH %']:
        GOH.append(min(i, 100))
    engine_goh['GOH'] = GOH

    # Alternator OH
    alternator_oh = [['GEN01', last['Last Overhaul GEN01'],],
                    ['GEN02', last['Last Overhaul GEN02'],],
                    ['GEN03', last['Last Overhaul GEN03'],],
                    ['GEN04', last['Last Overhaul GEN04'],]]
    alternator_oh = pd.DataFrame(alternator_oh, columns=['Sumba', 'Last Overhaul (Date)',])

    last_oh = []
    for i in alternator_oh['Last Overhaul (Date)']:
        last_oh.append(i+relativedelta(months=60))
    alternator_oh['Next Overhaul (Date)'] = last_oh

    alternator_oh['Todays Date'] = date.today()
    alternator_oh['Todays Date'] = alternator_oh['Todays Date'].astype('datetime64[ns]')
    alternator_oh['Time to Go'] = (alternator_oh['Next Overhaul (Date)'] - alternator_oh['Todays Date']).dt.days.astype('int64')
    alternator_oh['Five Years in Days'] = 1827
    alternator_oh['Utilization %'] = round((alternator_oh['Five Years in Days']-alternator_oh['Time to Go'])/ \
                                    alternator_oh['Five Years in Days']*100,0)

    util = []
    for i in alternator_oh['Utilization %']:
        util.append(min(i, 100))
    alternator_oh['Utilization'] = util
    alternator_oh['Max'] = 100

    #Critical Spare
    conveyor = last['Conveyor']
    df_conveyor = pd.DataFrame({'names' : ['progress',' '],
                    'values' :  [conveyor, 100 - conveyor]})

    engine = last['Engine']
    df_engine = pd.DataFrame({'names' : ['progress',' '],
                    'values' :  [engine, 100 - engine]})
    
    def shape_cu(value):
        shape=''
        if value == 'Green': shape="img/green check.png"
        elif value == 'Yellow': shape="img/yellow cross.png"
        elif value == 'Red': shape="img/red cross.png"
        elif value == 'Gray': shape="img/gray blank.png"
        return(shape)

    def shape_cs(value):
        shape=''
        if value<=(50) : shape="img/red cross.png"
        elif value<=(80) : shape="img/yellow cross.png"
        elif value<=100 : shape="img/green check.png"
        else : shape="img/gray blank.png"
        return(shape)

    def col_day_delta(value):
        col=''
        if value<=180 : col="img/red calendar.png"
        elif value<=365 : col="img/yellow calendar.png"
        else: col="img/green calendar.png"
        return(col)

    t_sumba = last['Next Docking Intermediate Survey (IS)']
    ConMakerShape_sumba = shape_cu(last['Conveyor-Maker Check-Up'])
    CraneMakerShape_sumba = shape_cu(last['Crane-Maker Check-Up'])
    EngMakerShape_sumba = shape_cu(last['Engine-Maker Check-Up'])
    OthersShape_sumba = shape_cu(last['Others'])
    NextDockingCol_sumba = col_day_delta((last['Next Docking Intermediate Survey (IS)'] - pd.Timestamp.today()).days)
    ConSpareShape_sumba = shape_cu(last['Conveyor Belt Spares'])
    CraneSpareShape_sumba = shape_cu(last['Crane-Wire Rope Spares']) 
    CraneShape_sumba = shape_cs(last['Crane'])
    ConShape_sumba = shape_cs(last['Conveyor'])
    EngShape_sumba = shape_cs(last['Engine'])

    #Chart
    sumba_hou_c1 = plot_hou1_s(hou1)
    sumba_hou_c2= plot_hou2_s(hou2)
    sumba_toh = plot_toh(engine_goh['TOH'])
    sumba_goh = plot_goh(engine_goh['GOH'])
    sumba_crane_1 = plot_wr1_s(wire_rope_c1['Utilization'], wire_rope_c1)
    sumba_crane_2 = plot_wr2_s(wire_rope_c2['Utilization'], wire_rope_c2)
    sumba_au = plot_aoh(alternator_oh['Utilization'], 'Sumba', alternator_oh)
    sumba_cs_conveyor = plot_cs(df_conveyor)
    sumba_cs_engine = plot_cs(df_engine)

    #Save Img
    pio.write_image(sumba_hou_c1, 'img/hou_c1 Bulk Sumba.png')
    pio.write_image(sumba_hou_c2, 'img/hou_c2 Bulk Sumba.png')
    pio.write_image(sumba_toh, 'img/toh Bulk Sumba.png')
    pio.write_image(sumba_goh, 'img/goh Bulk Sumba.png')
    pio.write_image(sumba_crane_1, 'img/crane_1 Bulk Sumba.png')
    pio.write_image(sumba_crane_2, 'img/crane_2 Bulk Sumba.png')
    pio.write_image(sumba_au, 'img/au Bulk Sumba.png')
    pio.write_image(sumba_cs_conveyor, 'img/cs_conveyor Bulk Sumba.png')
    pio.write_image(sumba_cs_engine, 'img/cs_engine Bulk Sumba.png')

    ###################### Java ######################
    df_java = pd.DataFrame(data['Bulk Java'])
    df_java = df_java.replace(r'^\s*$', np.nan, regex=True)

    df_java[['Last Overhaul GEN01', 'Last Overhaul GEN02', 'Last Overhaul GEN03', 'Last Overhaul GEN04']] = \
    df_java[['Last Overhaul GEN01', 'Last Overhaul GEN02', 'Last Overhaul GEN03', 'Last Overhaul GEN04']].fillna(date.today())
    df_java['Last Overhaul GEN01'] = pd.to_datetime(df_java['Last Overhaul GEN01'], dayfirst = True)
    df_java['Last Overhaul GEN02'] = pd.to_datetime(df_java['Last Overhaul GEN02'], dayfirst = True)
    df_java['Last Overhaul GEN03'] = pd.to_datetime(df_java['Last Overhaul GEN03'], dayfirst = True)
    df_java['Last Overhaul GEN04'] = pd.to_datetime(df_java['Last Overhaul GEN04'], dayfirst = True)
    df_java['Next Docking Intermediate Survey (IS)'] = pd.to_datetime(df_java['Next Docking Intermediate Survey (IS)'], dayfirst = True)

    # Hydrolic Oil Utilization
    last = df_java.iloc[-1]
    hou1 = (last['CRANE 1']-last['Last Rpl RH CR01'])/6000
    hou1 = round(hou1*100,0)
    hou1 = min(hou1, 100)

    hou2 = (last['CRANE 2']-last['Last Rpl RH CR02'])/6000
    hou2 = round(hou2*100,0)
    hou2 = min(hou2, 100)

    # Wire Rope
    wire_rope_c1 = [['HW RHOL', last['HW RHOL CR01'], last['CRANE 1']],
                    ['HW LHOL', last['HW LHOL CR01'], last['CRANE 1']],
                    ['CW RHOL', last['CW RHOL CR01'], last['CRANE 1']],
                    ['CW LHOL', last['CW LHOL CR01'], last['CRANE 1']]]
    wire_rope_c1 = pd.DataFrame(wire_rope_c1, columns=['Crane 1', 'Last Replacement', 'Current RH'])
    wire_rope_c1['Utilization RH'] = wire_rope_c1['Current RH']-wire_rope_c1['Last Replacement']
    wire_rope_c1['Utilization %'] = round(wire_rope_c1['Utilization RH']/4000*100,0)
    util1 = []
    for i in wire_rope_c1['Utilization %']:
        util1.append(min(i, 100))
    wire_rope_c1['Utilization'] = util1
    wire_rope_c1['Max'] = 100

    wire_rope_c2 = [['HW RHOL', last['HW RHOL CR02'], last['CRANE 2']],
                    ['HW LHOL', last['HW LHOL CR02'], last['CRANE 2']],
                    ['CW RHOL', last['CW RHOL CR02'], last['CRANE 2']],
                    ['CW LHOL', last['CW LHOL CR02'], last['CRANE 2']]]
    wire_rope_c2 = pd.DataFrame(wire_rope_c2, columns=['Crane 2', 'Last Replacement', 'Current RH'])
    wire_rope_c2['Utilization RH'] = wire_rope_c2['Current RH']-wire_rope_c2['Last Replacement']
    wire_rope_c2['Utilization %'] = round(wire_rope_c2['Utilization RH']/4000*100,0)
    util2 = []
    for i in wire_rope_c2['Utilization %']:
        util2.append(min(i, 100))
    wire_rope_c2['Utilization'] = util2
    wire_rope_c2['Max'] = 100

    # Engine GOH
    engine_goh = [['DG1', last['DG1'], last['Last TOH DG1'], last['Last GOH DG1']],
                ['DG2', last['DG2'], last['Last TOH DG2'], last['Last GOH DG2']],
                ['DG3', last['DG3'], last['Last TOH DG3'], last['Last GOH DG3']],
                ['EDG', last['EDG'], last['Last TOH EDG'], last['Last GOH EDG']]]
    engine_goh = pd.DataFrame(engine_goh, columns=['Java', 'Current HR', 'Last TOH', 'Last GOH'])
    engine_goh['Last TOH'] = engine_goh['Last TOH'].fillna(0)
    engine_goh['Utilization TOH'] = engine_goh['Current HR']-engine_goh['Last TOH']
    engine_goh['Utilization GOH'] = engine_goh['Current HR']-engine_goh['Last GOH']
    engine_goh['Hours till TOH'] = 12000-engine_goh['Utilization TOH']
    engine_goh['Hours till GOH'] = 24000-engine_goh['Utilization GOH']
    engine_goh['TOH %'] = round((1-engine_goh['Hours till TOH']/12000)*100,0)
    engine_goh['GOH %'] = round((1-engine_goh['Hours till GOH']/24000)*100,0)

    TOH = []
    for i in engine_goh['TOH %']:
        TOH.append(min(i, 100))
    engine_goh['TOH'] = TOH

    GOH = []
    for i in engine_goh['GOH %']:
        GOH.append(min(i, 100))
    engine_goh['GOH'] = GOH

    # Conveyor Belt Replacement
    conveyor_br = [['BC1', last['BC1'], last['Last RPL HR BC1']],
                ['BC2', last['BC2'], last['Last RPL HR BC2']],
                ['BC3', last['BC3'], last['Last RPL HR BC3']],
                ['BC4', last['BC4'], last['Last RPL HR BC4']],
                ['BC5', last['BC5'], last['Last RPL HR BC5']],
                ['SHL1', last['SHL1'], last['Last RPL HR SHL1']],
                ['SHL2', last['SHL2'], last['Last RPL HR SHL2']]]
    conveyor_br = pd.DataFrame(conveyor_br, columns=['Java', 'Current RH', 'Last Rpl HR'])
    conveyor_br['Utilization (hrs)'] = conveyor_br['Current RH']-conveyor_br['Last Rpl HR']
    conveyor_br['Utilization %'] = [round(conveyor_br['Utilization (hrs)'][0]/8000*100,0),
                                round(conveyor_br['Utilization (hrs)'][1]/8000*100,0),
                                round(conveyor_br['Utilization (hrs)'][2]/5000*100,0),
                                round(conveyor_br['Utilization (hrs)'][3]/5000*100,0),
                                round(conveyor_br['Utilization (hrs)'][4]/5000*100,0),
                                round(conveyor_br['Utilization (hrs)'][5]/5000*100,0),
                                round(conveyor_br['Utilization (hrs)'][6]/5000*100,0),]

    util = []
    for i in conveyor_br['Utilization %']:
        util.append(min(i, 100))
    conveyor_br['Utilization'] = util
    conveyor_br['Max'] = 100

    # Alternator OH
    alternator_oh = [['GEN01', last['Last Overhaul GEN01'],],
                    ['GEN02', last['Last Overhaul GEN02'],],
                    ['GEN03', last['Last Overhaul GEN03'],],
                    ['GEN04', last['Last Overhaul GEN04'],]]
    alternator_oh = pd.DataFrame(alternator_oh, columns=['Java', 'Last Overhaul (Date)',])

    last_oh = []
    for i in alternator_oh['Last Overhaul (Date)']:
        last_oh.append(i+relativedelta(months=60))
    alternator_oh['Next Overhaul (Date)'] = last_oh

    alternator_oh['Todays Date'] = date.today()
    alternator_oh['Todays Date'] = alternator_oh['Todays Date'].astype('datetime64[ns]')
    alternator_oh['Time to Go'] = (alternator_oh['Next Overhaul (Date)'] - alternator_oh['Todays Date']).dt.days.astype('int64')
    alternator_oh['Five Years in Days'] = 1827
    alternator_oh['Utilization %'] = round((alternator_oh['Five Years in Days']-alternator_oh['Time to Go'])/ \
                                    alternator_oh['Five Years in Days']*100,0)

    util = []
    for i in alternator_oh['Utilization %']:
        util.append(min(i, 100))
    alternator_oh['Utilization'] = util
    alternator_oh['Max'] = 100

    #Critical Spare
    crane = last['Crane']
    df_crane = pd.DataFrame({'names' : ['progress',' '],
                    'values' :  [crane, 100 - crane]})

    conveyor = last['Conveyor']
    df_conveyor = pd.DataFrame({'names' : ['progress',' '],
                    'values' :  [conveyor, 100 - conveyor]})

    engine = last['Engine']
    df_engine = pd.DataFrame({'names' : ['progress',' '],
                    'values' :  [engine, 100 - engine]})
    
    def shape_cu(value):
        shape=''
        if value == 'Green': shape="img/green check.png"
        elif value == 'Yellow': shape="img/yellow cross.png"
        elif value == 'Red': shape="img/red cross.png"
        elif value == 'Gray': shape="img/gray blank.png"
        return(shape)

    def shape_cs(value):
        shape=''
        if value<=(50) : shape="img/red cross.png"
        elif value<=(80) : shape="img/yellow cross.png"
        elif value<=100 : shape="img/green check.png"
        else : shape="img/gray blank.png"
        return(shape)

    def col_day_delta(value):
        col=''
        if value<=180 : col="img/red calendar.png"
        elif value<=365 : col="img/yellow calendar.png"
        else: col="img/green calendar.png"
        return(col)

    t_java = last['Next Docking Intermediate Survey (IS)']
    ConMakerShape_java = shape_cu(last['Conveyor-Maker Check-Up'])
    CraneMakerShape_java = shape_cu(last['Crane-Maker Check-Up'])
    EngMakerShape_java = shape_cu(last['Engine-Maker Check-Up'])
    OthersShape_java = shape_cu(last['Others'])
    NextDockingCol_java = col_day_delta((last['Next Docking Intermediate Survey (IS)'] - pd.Timestamp.today()).days)
    ConSpareShape_java = shape_cu(last['Conveyor Belt Spares'])
    CraneSpareShape_java = shape_cu(last['Crane-Wire Rope Spares']) 
    CraneShape_java = shape_cs(last['Crane'])
    ConShape_java = shape_cs(last['Conveyor'])
    EngShape_java = shape_cs(last['Engine'])

    #Chart
    java_hou_c1 = plot_hou1(hou1)
    java_hou_c2= plot_hou2(hou2)
    java_toh = plot_toh(engine_goh['TOH'])
    java_goh = plot_goh(engine_goh['GOH'])
    java_crane_1 = plot_wr1(wire_rope_c1['Utilization'], wire_rope_c1)
    java_crane_2 = plot_wr2(wire_rope_c2['Utilization'], wire_rope_c2)
    java_cbu = plot_cbr(conveyor_br['Utilization'], 'Java', conveyor_br)
    java_au = plot_aoh(alternator_oh['Utilization'], 'Java', alternator_oh)
    java_cs_crane = plot_cs(df_crane)
    java_cs_conveyor = plot_cs(df_conveyor)
    java_cs_engine = plot_cs(df_engine)

    #Save Img
    pio.write_image(java_hou_c1, 'img/hou_c1 Bulk Java.png')
    pio.write_image(java_hou_c2, 'img/hou_c2 Bulk Java.png')
    pio.write_image(java_toh, 'img/toh Bulk Java.png')
    pio.write_image(java_goh, 'img/goh Bulk Java.png')
    pio.write_image(java_crane_1, 'img/crane_1 Bulk Java.png')
    pio.write_image(java_crane_2, 'img/crane_2 Bulk Java.png')
    pio.write_image(java_cbu, 'img/cbu Bulk Java.png')
    pio.write_image(java_au, 'img/au Bulk Java.png')
    pio.write_image(java_cs_crane, 'img/cs_crane Bulk Java.png')
    pio.write_image(java_cs_conveyor, 'img/cs_conveyor Bulk Java.png')
    pio.write_image(java_cs_engine, 'img/cs_engine Bulk Java.png')

    ###################### Borneo ######################
    df_borneo = pd.DataFrame(data['Bulk Borneo'])
    df_borneo = df_borneo.replace(r'^\s*$', np.nan, regex=True)

    df_borneo[['Last Overhaul GEN01', 'Last Overhaul GEN02', 'Last Overhaul GEN03', 'Last Overhaul GEN04']] = \
    df_borneo[['Last Overhaul GEN01', 'Last Overhaul GEN02', 'Last Overhaul GEN03', 'Last Overhaul GEN04']].fillna(date.today())
    df_borneo['Last Overhaul GEN01'] = pd.to_datetime(df_borneo['Last Overhaul GEN01'], dayfirst = True)
    df_borneo['Last Overhaul GEN02'] = pd.to_datetime(df_borneo['Last Overhaul GEN02'], dayfirst = True)
    df_borneo['Last Overhaul GEN03'] = pd.to_datetime(df_borneo['Last Overhaul GEN03'], dayfirst = True)
    df_borneo['Last Overhaul GEN04'] = pd.to_datetime(df_borneo['Last Overhaul GEN04'], dayfirst = True)
    df_borneo['Next Docking Intermediate Survey (IS)'] = pd.to_datetime(df_borneo['Next Docking Intermediate Survey (IS)'], dayfirst = True)

    # Hydrolic Oil Utilization
    last = df_borneo.iloc[-1]
    hou1 = (last['CRANE 1']-last['Last Rpl RH CR01'])/6000
    hou1 = round(hou1*100,0)
    hou1 = min(hou1, 100)

    hou2 = (last['CRANE 2']-last['Last Rpl RH CR02'])/6000
    hou2 = round(hou2*100,0)
    hou2 = min(hou2, 100)

    # Wire Rope
    wire_rope_c1 = [['HW RHOL', last['HW RHOL CR01'], last['CRANE 1']],
                    ['HW LHOL', last['HW LHOL CR01'], last['CRANE 1']],
                    ['CW RHOL', last['CW RHOL CR01'], last['CRANE 1']],
                    ['CW LHOL', last['CW LHOL CR01'], last['CRANE 1']]]
    wire_rope_c1 = pd.DataFrame(wire_rope_c1, columns=['Crane 1', 'Last Replacement', 'Current RH'])
    wire_rope_c1['Utilization RH'] = wire_rope_c1['Current RH']-wire_rope_c1['Last Replacement']
    wire_rope_c1['Utilization %'] = round(wire_rope_c1['Utilization RH']/4000*100,0)
    util1 = []
    for i in wire_rope_c1['Utilization %']:
        util1.append(min(i, 100))
    wire_rope_c1['Utilization'] = util1
    wire_rope_c1['Max'] = 100

    wire_rope_c2 = [['HW RHOL', last['HW RHOL CR02'], last['CRANE 2']],
                    ['HW LHOL', last['HW LHOL CR02'], last['CRANE 2']],
                    ['CW RHOL', last['CW RHOL CR02'], last['CRANE 2']],
                    ['CW LHOL', last['CW LHOL CR02'], last['CRANE 2']]]
    wire_rope_c2 = pd.DataFrame(wire_rope_c2, columns=['Crane 2', 'Last Replacement', 'Current RH'])
    wire_rope_c2['Utilization RH'] = wire_rope_c2['Current RH']-wire_rope_c2['Last Replacement']
    wire_rope_c2['Utilization %'] = round(wire_rope_c2['Utilization RH']/4000*100,0)
    util2 = []
    for i in wire_rope_c2['Utilization %']:
        util2.append(min(i, 100))
    wire_rope_c2['Utilization'] = util2
    wire_rope_c2['Max'] = 100

    # Engine GOH
    engine_goh = [['DG1', last['DG1'], last['Last TOH DG1'], last['Last GOH DG1']],
                ['DG2', last['DG2'], last['Last TOH DG2'], last['Last GOH DG2']],
                ['DG3', last['DG3'], last['Last TOH DG3'], last['Last GOH DG3']],
                ['EDG', last['EDG'], last['Last TOH EDG'], last['Last GOH EDG']]]
    engine_goh = pd.DataFrame(engine_goh, columns=['Borneo', 'Current HR', 'Last TOH', 'Last GOH'])
    engine_goh['Last TOH'] = engine_goh['Last TOH'].fillna(0)
    engine_goh['Utilization TOH'] = engine_goh['Current HR']-engine_goh['Last TOH']
    engine_goh['Utilization GOH'] = engine_goh['Current HR']-engine_goh['Last GOH']
    engine_goh['Hours till TOH'] = 12000-engine_goh['Utilization TOH']
    engine_goh['Hours till GOH'] = 24000-engine_goh['Utilization GOH']
    engine_goh['TOH %'] = round((1-engine_goh['Hours till TOH']/12000)*100,0)
    engine_goh['GOH %'] = round((1-engine_goh['Hours till GOH']/24000)*100,0)

    TOH = []
    for i in engine_goh['TOH %']:
        TOH.append(min(i, 100))
    engine_goh['TOH'] = TOH

    GOH = []
    for i in engine_goh['GOH %']:
        GOH.append(min(i, 100))
    engine_goh['GOH'] = GOH

    # Conveyor Belt Replacement
    conveyor_br = [['BC1', last['BC1'], last['Last RPL HR BC1']],
                ['BC2', last['BC2'], last['Last RPL HR BC2']],
                ['BC3', last['BC3'], last['Last RPL HR BC3']],
                ['BC4', last['BC4'], last['Last RPL HR BC4']],
                ['BC5', last['BC5'], last['Last RPL HR BC5']],
                ['SHL1', last['SHL1'], last['Last RPL HR SHL1']],
                ['SHL2', last['SHL2'], last['Last RPL HR SHL2']]]
    conveyor_br = pd.DataFrame(conveyor_br, columns=['Borneo', 'Current RH', 'Last Rpl HR'])
    conveyor_br['Utilization (hrs)'] = conveyor_br['Current RH']-conveyor_br['Last Rpl HR']
    conveyor_br['Utilization %'] = [round(conveyor_br['Utilization (hrs)'][0]/6400*100,0),
                                round(conveyor_br['Utilization (hrs)'][1]/6400*100,0),
                                round(conveyor_br['Utilization (hrs)'][2]/4000*100,0),
                                round(conveyor_br['Utilization (hrs)'][3]/4000*100,0),
                                round(conveyor_br['Utilization (hrs)'][4]/4000*100,0),
                                round(conveyor_br['Utilization (hrs)'][5]/4000*100,0),
                                round(conveyor_br['Utilization (hrs)'][6]/4000*100,0),]

    util = []
    for i in conveyor_br['Utilization %']:
        util.append(min(i, 100))
    conveyor_br['Utilization'] = util
    conveyor_br['Max'] = 100

    # Alternator OH
    alternator_oh = [['GEN01', last['Last Overhaul GEN01'],],
                    ['GEN02', last['Last Overhaul GEN02'],],
                    ['GEN03', last['Last Overhaul GEN03'],],
                    ['GEN04', last['Last Overhaul GEN04'],]]
    alternator_oh = pd.DataFrame(alternator_oh, columns=['Borneo', 'Last Overhaul (Date)',])

    last_oh = []
    for i in alternator_oh['Last Overhaul (Date)']:
        last_oh.append(i+relativedelta(months=60))
    alternator_oh['Next Overhaul (Date)'] = last_oh

    alternator_oh['Todays Date'] = date.today()
    alternator_oh['Todays Date'] = alternator_oh['Todays Date'].astype('datetime64[ns]')
    alternator_oh['Time to Go'] = (alternator_oh['Next Overhaul (Date)'] - alternator_oh['Todays Date']).dt.days.astype('int64')
    alternator_oh['Five Years in Days'] = 1827
    alternator_oh['Utilization %'] = round((alternator_oh['Five Years in Days']-alternator_oh['Time to Go'])/ \
                                    alternator_oh['Five Years in Days']*100,0)

    util = []
    for i in alternator_oh['Utilization %']:
        util.append(min(i, 100))
    alternator_oh['Utilization'] = util
    alternator_oh['Max'] = 100

    #Critical Spare
    crane = last['Crane']
    df_crane = pd.DataFrame({'names' : ['progress',' '],
                    'values' :  [crane, 100 - crane]})

    conveyor = last['Conveyor']
    df_conveyor = pd.DataFrame({'names' : ['progress',' '],
                    'values' :  [conveyor, 100 - conveyor]})

    engine = last['Engine']
    df_engine = pd.DataFrame({'names' : ['progress',' '],
                    'values' :  [engine, 100 - engine]})
    
    def shape_cu(value):
        shape=''
        if value == 'Green': shape="img/green check.png"
        elif value == 'Yellow': shape="img/yellow cross.png"
        elif value == 'Red': shape="img/red cross.png"
        elif value == 'Gray': shape="img/gray blank.png"
        return(shape)

    def shape_cs(value):
        shape=''
        if value<=(50) : shape="img/red cross.png"
        elif value<=(80) : shape="img/yellow cross.png"
        elif value<=100 : shape="img/green check.png"
        else : shape="img/gray blank.png"
        return(shape)

    def col_day_delta(value):
        col=''
        if value<=180 : col="img/red calendar.png"
        elif value<=365 : col="img/yellow calendar.png"
        else: col="img/green calendar.png"
        return(col)

    t_borneo = last['Next Docking Intermediate Survey (IS)']
    ConMakerShape_borneo = shape_cu(last['Conveyor-Maker Check-Up'])
    CraneMakerShape_borneo = shape_cu(last['Crane-Maker Check-Up'])
    EngMakerShape_borneo = shape_cu(last['Engine-Maker Check-Up'])
    OthersShape_borneo = shape_cu(last['Others'])
    NextDockingCol_borneo = col_day_delta((last['Next Docking Intermediate Survey (IS)'] - pd.Timestamp.today()).days)
    ConSpareShape_borneo = shape_cu(last['Conveyor Belt Spares'])
    CraneSpareShape_borneo = shape_cu(last['Crane-Wire Rope Spares']) 
    CraneShape_borneo = shape_cs(last['Crane'])
    ConShape_borneo = shape_cs(last['Conveyor'])
    EngShape_borneo = shape_cs(last['Engine'])

    #Chart
    borneo_hou_c1 = plot_hou1(hou1)
    borneo_hou_c2= plot_hou2(hou2)
    borneo_toh = plot_toh(engine_goh['TOH'])
    borneo_goh = plot_goh(engine_goh['GOH'])
    borneo_crane_1 = plot_wr1(wire_rope_c1['Utilization'], wire_rope_c1)
    borneo_crane_2 = plot_wr2(wire_rope_c2['Utilization'], wire_rope_c2)
    borneo_cbu = plot_cbr(conveyor_br['Utilization'], 'Borneo', conveyor_br)
    borneo_au = plot_aoh(alternator_oh['Utilization'], 'Borneo', alternator_oh)
    borneo_cs_crane = plot_cs(df_crane)
    borneo_cs_conveyor = plot_cs(df_conveyor)
    borneo_cs_engine = plot_cs(df_engine)

    #Save Img
    pio.write_image(borneo_hou_c1, 'img/hou_c1 Bulk Borneo.png')
    pio.write_image(borneo_hou_c2, 'img/hou_c2 Bulk Borneo.png')
    pio.write_image(borneo_toh, 'img/toh Bulk Borneo.png')
    pio.write_image(borneo_goh, 'img/goh Bulk Borneo.png')
    pio.write_image(borneo_crane_1, 'img/crane_1 Bulk Borneo.png')
    pio.write_image(borneo_crane_2, 'img/crane_2 Bulk Borneo.png')
    pio.write_image(borneo_cbu, 'img/cbu Bulk Borneo.png')
    pio.write_image(borneo_au, 'img/au Bulk Borneo.png')
    pio.write_image(borneo_cs_crane, 'img/cs_crane Bulk Borneo.png')
    pio.write_image(borneo_cs_conveyor, 'img/cs_conveyor Bulk Borneo.png')
    pio.write_image(borneo_cs_engine, 'img/cs_engine Bulk Borneo.png')

    ###################### Celebes ######################
    df_celebes = pd.DataFrame(data['Bulk Celebes'])
    df_celebes = df_celebes.replace(r'^\s*$', np.nan, regex=True)

    df_celebes[['Last Overhaul GEN01', 'Last Overhaul GEN02', 'Last Overhaul GEN03', 'Last Overhaul GEN04']] = \
    df_celebes[['Last Overhaul GEN01', 'Last Overhaul GEN02', 'Last Overhaul GEN03', 'Last Overhaul GEN04']].fillna(date.today())
    df_celebes['Last Overhaul GEN01'] = pd.to_datetime(df_celebes['Last Overhaul GEN01'], dayfirst = True)
    df_celebes['Last Overhaul GEN02'] = pd.to_datetime(df_celebes['Last Overhaul GEN02'], dayfirst = True)
    df_celebes['Last Overhaul GEN03'] = pd.to_datetime(df_celebes['Last Overhaul GEN03'], dayfirst = True)
    df_celebes['Last Overhaul GEN04'] = pd.to_datetime(df_celebes['Last Overhaul GEN04'], dayfirst = True)
    df_celebes['Next Docking Intermediate Survey (IS)'] = pd.to_datetime(df_celebes['Next Docking Intermediate Survey (IS)'], dayfirst = True)

    # Hydrolic Oil Utilization
    last = df_celebes.iloc[-1]
    hou1 = (last['CRANE 1']-last['Last Rpl RH CR01'])/6000
    hou1 = round(hou1*100,0)
    hou1 = min(hou1, 100)

    hou2 = (last['CRANE 2']-last['Last Rpl RH CR02'])/6000
    hou2 = round(hou2*100,0)
    hou2 = min(hou2, 100)

    # Wire Rope
    wire_rope_c1 = [['HW RHOL', last['HW RHOL CR01'], last['CRANE 1']],
                    ['HW LHOL', last['HW LHOL CR01'], last['CRANE 1']],
                    ['CW RHOL', last['CW RHOL CR01'], last['CRANE 1']],
                    ['CW LHOL', last['CW LHOL CR01'], last['CRANE 1']],
                    ['EW RHOL', last['EW RHOL CR01'], last['CRANE 1']]]
    wire_rope_c1 = pd.DataFrame(wire_rope_c1, columns=['Crane 1', 'Last Replacement', 'Current RH'])
    wire_rope_c1['Utilization RH'] = wire_rope_c1['Current RH']-wire_rope_c1['Last Replacement']
    wire_rope_c1['Utilization %'] = [round(wire_rope_c1['Utilization RH'][0]/4000*100,0),
                                    round(wire_rope_c1['Utilization RH'][1]/4000*100,0),
                                    round(wire_rope_c1['Utilization RH'][2]/4000*100,0),
                                    round(wire_rope_c1['Utilization RH'][3]/4000*100,0),
                                    round(wire_rope_c1['Utilization RH'][4]/3500*100,0),]
    util1 = []
    for i in wire_rope_c1['Utilization %']:
        util1.append(min(i, 100))
    wire_rope_c1['Utilization'] = util1
    wire_rope_c1['Max'] = 100


    wire_rope_c2 = [['HW RHOL', last['HW RHOL CR02'], last['CRANE 2']],
                    ['HW LHOL', last['HW LHOL CR02'], last['CRANE 2']],
                    ['CW RHOL', last['CW RHOL CR02'], last['CRANE 2']],
                    ['CW LHOL', last['CW LHOL CR02'], last['CRANE 2']],
                    ['EW RHOL', last['EW RHOL CR02'], last['CRANE 2']]]
    wire_rope_c2 = pd.DataFrame(wire_rope_c2, columns=['Crane 2', 'Last Replacement', 'Current RH'])
    wire_rope_c2['Utilization RH'] = wire_rope_c2['Current RH']-wire_rope_c2['Last Replacement']
    wire_rope_c2['Utilization %'] = [round(wire_rope_c2['Utilization RH'][0]/4000*100,0),
                                    round(wire_rope_c2['Utilization RH'][1]/4000*100,0),
                                    round(wire_rope_c2['Utilization RH'][2]/4000*100,0),
                                    round(wire_rope_c2['Utilization RH'][3]/4000*100,0),
                                    round(wire_rope_c2['Utilization RH'][4]/3500*100,0),]
    util2 = []
    for i in wire_rope_c2['Utilization %']:
        util2.append(min(i, 100))
    wire_rope_c2['Utilization'] = util2
    wire_rope_c2['Max'] = 100

    # Engine GOH
    engine_goh = [['DG1', last['DG1'], last['Last TOH DG1'], last['Last GOH DG1']],
                ['DG2', last['DG2'], last['Last TOH DG2'], last['Last GOH DG2']],
                ['DG3', last['DG3'], last['Last TOH DG3'], last['Last GOH DG3']],
                ['EDG', last['EDG'], last['Last TOH EDG'], last['Last GOH EDG']]]
    engine_goh = pd.DataFrame(engine_goh, columns=['Celebes', 'Current HR', 'Last TOH', 'Last GOH'])
    engine_goh['Last TOH'] = engine_goh['Last TOH'].fillna(0)
    engine_goh['Utilization TOH'] = engine_goh['Current HR']-engine_goh['Last TOH']
    engine_goh['Utilization GOH'] = engine_goh['Current HR']-engine_goh['Last GOH']
    engine_goh['Hours till TOH'] = 12000-engine_goh['Utilization TOH']
    engine_goh['Hours till GOH'] = 24000-engine_goh['Utilization GOH']
    engine_goh['TOH %'] = round((1-engine_goh['Hours till TOH']/12000)*100,0)
    engine_goh['GOH %'] = round((1-engine_goh['Hours till GOH']/24000)*100,0)

    TOH = []
    for i in engine_goh['TOH %']:
        TOH.append(min(i, 100))
    engine_goh['TOH'] = TOH

    GOH = []
    for i in engine_goh['GOH %']:
        GOH.append(min(i, 100))
    engine_goh['GOH'] = GOH

    # Conveyor Belt Replacement
    conveyor_br = [['BC1', last['BC1'], last['Last RPL HR BC1']],
                ['BC2', last['BC2'], last['Last RPL HR BC2']],
                ['BC3', last['BC3'], last['Last RPL HR BC3']],
                ['BL', last['LC'], last['Last RPL HR BL']],
                ['BC4', last['BC4'], last['Last RPL HR BC4']],
                ['BC5', last['BC5'], last['Last RPL HR BC5']],
                ['BC6', last['BC6'], last['Last RPL HR BC6']],
                ['SHL1', last['SHL'], last['Last RPL HR SHL1']]]
    conveyor_br = pd.DataFrame(conveyor_br, columns=['Celebes', 'Current RH', 'Last Rpl HR'])
    conveyor_br['Utilization (hrs)'] = conveyor_br['Current RH']-conveyor_br['Last Rpl HR']
    conveyor_br['Utilization %'] = [round(conveyor_br['Utilization (hrs)'][0]/8000*100,0),
                                round(conveyor_br['Utilization (hrs)'][1]/8000*100,0),
                                round(conveyor_br['Utilization (hrs)'][2]/5000*100,0),
                                round(conveyor_br['Utilization (hrs)'][3]/5000*100,0),
                                round(conveyor_br['Utilization (hrs)'][4]/5000*100,0),
                                round(conveyor_br['Utilization (hrs)'][5]/5000*100,0),
                                round(conveyor_br['Utilization (hrs)'][6]/5000*100,0),
                                round(conveyor_br['Utilization (hrs)'][7]/5000*100,0),]

    util = []
    for i in conveyor_br['Utilization %']:
        util.append(min(i, 100))
    conveyor_br['Utilization'] = util
    conveyor_br['Max'] = 100


    # Alternator OH
    alternator_oh = [['GEN01', last['Last Overhaul GEN01'],],
                    ['GEN02', last['Last Overhaul GEN02'],],
                    ['GEN03', last['Last Overhaul GEN03'],],
                    ['GEN04', last['Last Overhaul GEN04'],]]
    alternator_oh = pd.DataFrame(alternator_oh, columns=['Celebes', 'Last Overhaul (Date)',])

    last_oh = []
    for i in alternator_oh['Last Overhaul (Date)']:
        last_oh.append(i+relativedelta(months=60))
    alternator_oh['Next Overhaul (Date)'] = last_oh

    alternator_oh['Todays Date'] = date.today()
    alternator_oh['Todays Date'] = alternator_oh['Todays Date'].astype('datetime64[ns]')
    alternator_oh['Time to Go'] = (alternator_oh['Next Overhaul (Date)'] - alternator_oh['Todays Date']).dt.days.astype('int64')
    alternator_oh['Five Years in Days'] = 1827
    alternator_oh['Utilization %'] = round((alternator_oh['Five Years in Days']-alternator_oh['Time to Go'])/ \
                                    alternator_oh['Five Years in Days']*100,0)

    util = []
    for i in alternator_oh['Utilization %']:
        util.append(min(i, 100))
    alternator_oh['Utilization'] = util
    alternator_oh['Max'] = 100

    #Critical Spare
    crane = last['Crane']
    df_crane = pd.DataFrame({'names' : ['progress',' '],
                    'values' :  [crane, 100 - crane]})

    conveyor = last['Conveyor']
    df_conveyor = pd.DataFrame({'names' : ['progress',' '],
                    'values' :  [conveyor, 100 - conveyor]})

    engine = last['Engine']
    df_engine = pd.DataFrame({'names' : ['progress',' '],
                    'values' :  [engine, 100 - engine]})
    
    def shape_cu(value):
        shape=''
        if value == 'Green': shape="img/green check.png"
        elif value == 'Yellow': shape="img/yellow cross.png"
        elif value == 'Red': shape="img/red cross.png"
        elif value == 'Gray': shape="img/gray blank.png"
        return(shape)

    def shape_cs(value):
        shape=''
        if value<=(50) : shape="img/red cross.png"
        elif value<=(80) : shape="img/yellow cross.png"
        elif value<=100 : shape="img/green check.png"
        else : shape="img/gray blank.png"
        return(shape)

    def col_day_delta(value):
        col=''
        if value<=180 : col="img/red calendar.png"
        elif value<=365 : col="img/yellow calendar.png"
        else: col="img/green calendar.png"
        return(col)

    t_celebes = last['Next Docking Intermediate Survey (IS)']
    ConMakerShape_celebes = shape_cu(last['Conveyor-Maker Check-Up'])
    CraneMakerShape_celebes = shape_cu(last['Crane-Maker Check-Up'])
    EngMakerShape_celebes = shape_cu(last['Engine-Maker Check-Up'])
    OthersShape_celebes = shape_cu(last['Others'])
    NextDockingCol_celebes = col_day_delta((last['Next Docking Intermediate Survey (IS)'] - pd.Timestamp.today()).days)
    ConSpareShape_celebes = shape_cu(last['Conveyor Belt Spares'])
    CraneSpareShape_celebes = shape_cu(last['Crane-Wire Rope Spares']) 
    CraneShape_celebes = shape_cs(last['Crane'])
    ConShape_celebes = shape_cs(last['Conveyor'])
    EngShape_celebes = shape_cs(last['Engine'])

    #Chart
    celebes_hou_c1 = plot_hou1(hou1)
    celebes_hou_c2= plot_hou2(hou2)
    celebes_toh = plot_toh(engine_goh['TOH'])
    celebes_goh = plot_goh(engine_goh['GOH'])
    celebes_crane_1 = plot_wr1(wire_rope_c1['Utilization'], wire_rope_c1)
    celebes_crane_2 = plot_wr2(wire_rope_c2['Utilization'], wire_rope_c2)
    celebes_cbu = plot_cbr(conveyor_br['Utilization'], 'Celebes', conveyor_br)
    celebes_au = plot_aoh(alternator_oh['Utilization'], 'Celebes', alternator_oh)
    celebes_cs_crane = plot_cs(df_crane)
    celebes_cs_conveyor = plot_cs(df_conveyor)
    celebes_cs_engine = plot_cs(df_engine)

    #Save Img
    pio.write_image(celebes_hou_c1, 'img/hou_c1 Bulk Celebes.png')
    pio.write_image(celebes_hou_c2, 'img/hou_c2 Bulk Celebes.png')
    pio.write_image(celebes_toh, 'img/toh Bulk Celebes.png')
    pio.write_image(celebes_goh, 'img/goh Bulk Celebes.png')
    pio.write_image(celebes_crane_1, 'img/crane_1 Bulk Celebes.png')
    pio.write_image(celebes_crane_2, 'img/crane_2 Bulk Celebes.png')
    pio.write_image(celebes_cbu, 'img/cbu Bulk Celebes.png')
    pio.write_image(celebes_au, 'img/au Bulk Celebes.png')
    pio.write_image(celebes_cs_crane, 'img/cs_crane Bulk Celebes.png')
    pio.write_image(celebes_cs_conveyor, 'img/cs_conveyor Bulk Celebes.png')
    pio.write_image(celebes_cs_engine, 'img/cs_engine Bulk Celebes.png')


    ###################### Natuna ######################
    df_natuna = pd.DataFrame(data['Bulk Natuna'])
    df_natuna = df_natuna.replace(r'^\s*$', np.nan, regex=True)

    df_natuna[['Last Overhaul GEN01', 'Last Overhaul GEN02', 'Last Overhaul GEN03', 'Last Overhaul GEN04']] = \
    df_natuna[['Last Overhaul GEN01', 'Last Overhaul GEN02', 'Last Overhaul GEN03', 'Last Overhaul GEN04']].fillna(date.today())
    df_natuna['Last Overhaul GEN01'] = pd.to_datetime(df_natuna['Last Overhaul GEN01'], dayfirst = True)
    df_natuna['Last Overhaul GEN02'] = pd.to_datetime(df_natuna['Last Overhaul GEN02'], dayfirst = True)
    df_natuna['Last Overhaul GEN03'] = pd.to_datetime(df_natuna['Last Overhaul GEN03'], dayfirst = True)
    df_natuna['Last Overhaul GEN04'] = pd.to_datetime(df_natuna['Last Overhaul GEN04'], dayfirst = True)
    df_natuna['Next Docking Intermediate Survey (IS)'] = pd.to_datetime(df_natuna['Next Docking Intermediate Survey (IS)'], dayfirst = True)


    # Hydrolic Oil Utilization
    last = df_natuna.iloc[-1]
    hou1 = (last['CRANE']-last['Last Rpl RH CR01'])/6000
    hou1 = round(hou1*100,0)
    hou1 = min(hou1, 100)

    # Wire Rope
    wire_rope_c1 = [['HW RHOL', last['HW RHOL CR01'], last['CRANE']],
                    ['HW LHOL', last['HW LHOL CR01'], last['CRANE']],
                    ['CW RHOL', last['CW RHOL CR01'], last['CRANE']],
                    ['CW LHOL', last['CW LHOL CR01'], last['CRANE']],
                    ['EW RHOL', last['EW RHOL CR01'], last['CRANE']],]
    wire_rope_c1 = pd.DataFrame(wire_rope_c1, columns=['Crane 1', 'Last Replacement', 'Current RH'])
    wire_rope_c1['Utilization RH'] = wire_rope_c1['Current RH']-wire_rope_c1['Last Replacement']
    wire_rope_c1['Utilization %'] = [round(wire_rope_c1['Utilization RH'][0]/4000*100,0),
                                    round(wire_rope_c1['Utilization RH'][1]/4000*100,0),
                                    round(wire_rope_c1['Utilization RH'][2]/4000*100,0),
                                    round(wire_rope_c1['Utilization RH'][3]/4000*100,0),
                                    round(wire_rope_c1['Utilization RH'][4]/2000*100,0),]
    util1 = []
    for i in wire_rope_c1['Utilization %']:
        util1.append(min(i, 100))
    wire_rope_c1['Utilization'] = util1
    wire_rope_c1['Max'] = 100


    # Engine GOH
    engine_goh = [['DG1', last['HARBOR GEN 1'], last['Last TOH DG1'], last['Last GOH DG1']],
                ['DG2', last['HARBOR GEN 2'], last['Last TOH DG2'], last['Last GOH DG2']],
                ['DG3', last['MAIN GENERATOR 3'], last['Last TOH DG3'], last['Last GOH DG3']],
                ['EDG', last['MAIN GENERATOR 4'], last['Last TOH EDG'], last['Last GOH EDG']]]
    engine_goh = pd.DataFrame(engine_goh, columns=['Natuna', 'Current HR', 'Last TOH', 'Last GOH'])
    engine_goh['Last TOH'] = engine_goh['Last TOH'].fillna(0)
    engine_goh['Utilization TOH'] = engine_goh['Current HR']-engine_goh['Last TOH']
    engine_goh['Utilization GOH'] = engine_goh['Current HR']-engine_goh['Last GOH']
    engine_goh['Hours till TOH'] = 12000-engine_goh['Utilization TOH']
    engine_goh['Hours till GOH'] = 24000-engine_goh['Utilization GOH']
    engine_goh['TOH %'] = round((1-engine_goh['Hours till TOH']/12000)*100,0)
    engine_goh['GOH %'] = round((1-engine_goh['Hours till GOH']/24000)*100,0)

    TOH = []
    for i in engine_goh['TOH %']:
        TOH.append(min(i, 100))
    engine_goh['TOH'] = TOH

    GOH = []
    for i in engine_goh['GOH %']:
        GOH.append(min(i, 100))
    engine_goh['GOH'] = GOH

    # Alternator OH
    alternator_oh = [['GEN01', last['Last Overhaul GEN01'],],
                    ['GEN02', last['Last Overhaul GEN02'],],
                    ['GEN03', last['Last Overhaul GEN03'],],
                    ['GEN04', last['Last Overhaul GEN04'],]]
    alternator_oh = pd.DataFrame(alternator_oh, columns=['Natuna', 'Last Overhaul (Date)',])

    last_oh = []
    for i in alternator_oh['Last Overhaul (Date)']:
        last_oh.append(i+relativedelta(months=60))
    alternator_oh['Next Overhaul (Date)'] = last_oh

    alternator_oh['Todays Date'] = date.today()
    alternator_oh['Todays Date'] = alternator_oh['Todays Date'].astype('datetime64[ns]')
    alternator_oh['Time to Go'] = (alternator_oh['Next Overhaul (Date)'] - alternator_oh['Todays Date']).dt.days.astype('int64')
    alternator_oh['Five Years in Days'] = 1827
    alternator_oh['Utilization %'] = round((alternator_oh['Five Years in Days']-alternator_oh['Time to Go'])/ \
                                    alternator_oh['Five Years in Days']*100,0)

    util = []
    for i in alternator_oh['Utilization %']:
        util.append(min(i, 100))
    alternator_oh['Utilization'] = util
    alternator_oh['Max'] = 100

    def shape_cu(value):
        shape=''
        if value == 'Green': shape="img/green check.png"
        elif value == 'Yellow': shape="img/yellow cross.png"
        elif value == 'Red': shape="img/red cross.png"
        elif value == 'Gray': shape="img/gray blank.png"
        return(shape)

    def shape_cs(value):
        shape=''
        if value<=(50) : shape="img/red cross.png"
        elif value<=(80) : shape="img/yellow cross.png"
        elif value<=100 : shape="img/green check.png"
        else : shape="img/gray blank.png"
        return(shape)

    def col_day_delta(value):
        col=''
        if value<=180 : col="img/red calendar.png"
        elif value<=365 : col="img/yellow calendar.png"
        else: col="img/green calendar.png"
        return(col)

    t_natuna = last['Next Docking Intermediate Survey (IS)']
    ConMakerShape_natuna = shape_cu(last['Conveyor-Maker Check-Up'])
    CraneMakerShape_natuna = shape_cu(last['Crane-Maker Check-Up'])
    EngMakerShape_natuna = shape_cu(last['Engine-Maker Check-Up'])
    OthersShape_natuna = shape_cu(last['Others'])
    NextDockingCol_natuna = col_day_delta((last['Next Docking Intermediate Survey (IS)'] - pd.Timestamp.today()).days)
    ConSpareShape_natuna = shape_cu(last['Conveyor Belt Spares'])
    CraneSpareShape_natuna = shape_cu(last['Crane-Wire Rope Spares']) 
    CraneShape_natuna = shape_cs(last['Crane'])
    ConShape_natuna = shape_cs(last['Conveyor'])
    EngShape_natuna = shape_cs(last['Engine'])

    #Chart
    natuna_hou_c1 = plot_hou1_n(hou1)
    natuna_toh = plot_toh(engine_goh['TOH'])
    natuna_goh = plot_goh(engine_goh['GOH'])
    natuna_crane_1 = plot_wr1_n(wire_rope_c1['Utilization'], wire_rope_c1)
    natuna_au = plot_aoh(alternator_oh['Utilization'], 'Natuna', alternator_oh)
    #Save Img
    pio.write_image(natuna_hou_c1, 'img/hou_c1 Bulk Natuna.png')
    pio.write_image(natuna_toh, 'img/toh Bulk Natuna.png')
    pio.write_image(natuna_goh, 'img/goh Bulk Natuna.png')
    pio.write_image(natuna_crane_1, 'img/crane_1 Bulk Natuna.png')
    pio.write_image(natuna_au, 'img/au Bulk Natuna.png')

    ###################### PPTX ######################
    def to_pptx(bytes_io):
        # Create presentation
        pptx = path + '//' + 'slide_master.pptx'
        prs = Presentation(pptx)

        # define slidelayouts 
        slide1 = prs.slides.add_slide(prs.slide_layouts[0])
        slide2 = prs.slides.add_slide(prs.slide_layouts[0])
        slide3 = prs.slides.add_slide(prs.slide_layouts[0])
        slide4 = prs.slides.add_slide(prs.slide_layouts[0])
        slide5 = prs.slides.add_slide(prs.slide_layouts[0])
        slide6 = prs.slides.add_slide(prs.slide_layouts[0])
        slide7 = prs.slides.add_slide(prs.slide_layouts[0])
        slide8 = prs.slides.add_slide(prs.slide_layouts[0])
        slide9 = prs.slides.add_slide(prs.slide_layouts[0])

        # slide1 Bulk Sumatra
        slide1.placeholders[10].text = 'EQUIPMENT DASHBOARD \n Bulk Sumatra'
        slide1.shapes.add_picture('img/hou_c1 Bulk Sumatra.png', Cm(0.75), Cm(3.15), Cm(4), Cm(4))
        slide1.shapes.add_picture('img/hou_c2 Bulk Sumatra.png', Cm(0.75), Cm(7.4), Cm(4), Cm(4))
        slide1.shapes.add_picture('img/crane_1 Bulk Sumatra.png', Cm(5), Cm(3.15), Cm(9.5), Cm(4))
        slide1.shapes.add_picture('img/crane_2 Bulk Sumatra.png', Cm(5), Cm(7.4), Cm(9.5), Cm(4))
        slide1.shapes.add_picture('img/toh Bulk Sumatra.png', Cm(14.75), Cm(3.15), Cm(11), Cm(4))
        slide1.shapes.add_picture('img/goh Bulk Sumatra.png', Cm(14.75), Cm(7.4), Cm(11), Cm(4))
        slide1.shapes.add_picture('img/cbu Bulk Sumatra.png', Cm(0.75), Cm(11.65), Cm(13.75), Cm(5.5))
        slide1.shapes.add_picture('img/au Bulk Sumatra.png', Cm(14.75), Cm(11.65), Cm(11), Cm(5.5))
        slide1.shapes.add_picture('img/cs_crane Bulk Sumatra.png', Cm(26.2), Cm(14.2), Cm(2), Cm(2.5))
        slide1.shapes.add_picture('img/cs_conveyor Bulk Sumatra.png', Cm(28.5), Cm(14.2), Cm(2), Cm(2.5))
        slide1.shapes.add_picture('img/cs_engine Bulk Sumatra.png', Cm(30.9), Cm(14.2), Cm(2), Cm(2.5))

        slide1.placeholders[12].text = 'Conveyor-Maker Check-Up'
        slide1.placeholders[13].text = 'Crane-Maker Check-Up'
        slide1.placeholders[14].text = 'Engine-Maker Check-Up'
        slide1.placeholders[15].text = 'Others-Shi ELD'
        slide1.placeholders[16].text = 'Next Docking \n Intermediate \n Survey (IS)'
        slide1.placeholders[17].text = 'Conveyor Belt Spares'
        slide1.placeholders[18].text = 'Crane-Wire Rope Spares'
        slide1.placeholders[19].text = 'Crane'
        slide1.placeholders[20].text = 'Conv.'
        slide1.placeholders[21].text = 'Eng.'

        cal = slide1.shapes.add_picture(NextDockingCol_sumatra, Cm(30), Cm(7.7))

        thn = slide1.shapes.add_textbox(Cm(30.6), Cm(8), 1, 0.5)
        thn = thn.text_frame.paragraphs[0]
        thn.text = t_sumatra.strftime("%Y")
        font = thn.font
        font.bold = True
        font.name = 'Arial Black'
        font.size = Pt(14)

        tgl = slide1.shapes.add_textbox(Cm(30.8), Cm(8.6), 1, 0.5)
        tgl = tgl.text_frame.paragraphs[0]
        tgl.text = t_sumatra.strftime("%d")
        font = tgl.font
        font.bold = True
        font.name = 'Arial Black'
        font.size = Pt(18)

        bln = slide1.shapes.add_textbox(Cm(30.8), Cm(9.35), 1, 0.5)
        bln = bln.text_frame.paragraphs[0]
        bln.text = t_sumatra.strftime("%b")
        font = bln.font
        font.name = 'Arial'
        font.size = Pt(14)

        slide1.shapes.add_picture(ConMakerShape_sumatra, Cm(32.1), Cm(3.3))
        slide1.shapes.add_picture(CraneMakerShape_sumatra, Cm(32.1), Cm(4.4))
        slide1.shapes.add_picture(EngMakerShape_sumatra, Cm(32.1), Cm(5.5))
        slide1.shapes.add_picture(OthersShape_sumatra, Cm(32.1), Cm(6.6))
        slide1.shapes.add_picture(ConSpareShape_sumatra, Cm(32.1), Cm(10.7))
        slide1.shapes.add_picture(CraneSpareShape_sumatra, Cm(32.1), Cm(11.8))
        slide1.shapes.add_picture(CraneShape_sumatra, Cm(27.4), Cm(12.9))
        slide1.shapes.add_picture(ConShape_sumatra, Cm(29.8), Cm(12.9))
        slide1.shapes.add_picture(EngShape_sumatra, Cm(32.1), Cm(12.9))

        # slide2 Bulk Derawan
        slide2.placeholders[10].text = 'EQUIPMENT DASHBOARD \n Bulk Derawan'
        slide2.shapes.add_picture('img/crane_1 Bulk Derawan.png', Cm(0.75), Cm(3.15), Cm(13.75), Cm(4))
        slide2.shapes.add_picture('img/crane_2 Bulk Derawan.png', Cm(0.75), Cm(7.4), Cm(13.75), Cm(4))
        slide2.shapes.add_picture('img/toh Bulk Derawan.png', Cm(14.75), Cm(3.15), Cm(11), Cm(4))
        slide2.shapes.add_picture('img/goh Bulk Derawan.png', Cm(14.75), Cm(7.4), Cm(11), Cm(4))
        slide2.shapes.add_picture('img/cbu Bulk Derawan.png', Cm(0.75), Cm(11.65), Cm(13.75), Cm(5.5))
        slide2.shapes.add_picture('img/au Bulk Derawan.png', Cm(14.75), Cm(11.65), Cm(11), Cm(5.5))
        slide2.shapes.add_picture('img/cs_conveyor Bulk Derawan.png', Cm(28.5), Cm(14.2), Cm(2), Cm(2.5))


        slide2.placeholders[12].text = 'Conveyor-Maker Check-Up'
        slide2.placeholders[13].text = 'Crane-Maker Check-Up'
        slide2.placeholders[14].text = 'Engine-Maker Check-Up'
        slide2.placeholders[15].text = 'Others-Shi ELD'
        slide2.placeholders[16].text = 'Next Docking \n Intermediate \n Survey (IS)'
        slide2.placeholders[17].text = 'Conveyor Belt Spares'
        slide2.placeholders[18].text = 'Crane-Wire Rope Spares'
        slide2.placeholders[19].text = 'Crane'
        slide2.placeholders[20].text = 'Conv.'
        slide2.placeholders[21].text = 'Eng.'

        cal = slide2.shapes.add_picture(NextDockingCol_derawan, Cm(30), Cm(7.7))

        thn = slide2.shapes.add_textbox(Cm(30.6), Cm(8), 1, 0.5)
        thn = thn.text_frame.paragraphs[0]
        thn.text = t_derawan.strftime("%Y")
        font = thn.font
        font.bold = True
        font.name = 'Arial Black'
        font.size = Pt(14)

        tgl = slide2.shapes.add_textbox(Cm(30.8), Cm(8.6), 1, 0.5)
        tgl = tgl.text_frame.paragraphs[0]
        tgl.text = t_derawan.strftime("%d")
        font = tgl.font
        font.bold = True
        font.name = 'Arial Black'
        font.size = Pt(18)

        bln = slide2.shapes.add_textbox(Cm(30.8), Cm(9.35), 1, 0.5)
        bln = bln.text_frame.paragraphs[0]
        bln.text = t_derawan.strftime("%b")
        font = bln.font
        font.name = 'Arial'
        font.size = Pt(14)

        slide2.shapes.add_picture(ConMakerShape_derawan, Cm(32.1), Cm(3.3))
        slide2.shapes.add_picture(CraneMakerShape_derawan, Cm(32.1), Cm(4.4))
        slide2.shapes.add_picture(EngMakerShape_derawan, Cm(32.1), Cm(5.5))
        slide2.shapes.add_picture(OthersShape_derawan, Cm(32.1), Cm(6.6))
        slide2.shapes.add_picture(ConSpareShape_derawan, Cm(32.1), Cm(10.7))
        slide2.shapes.add_picture(CraneSpareShape_derawan, Cm(32.1), Cm(11.8))
        slide2.shapes.add_picture(CraneShape_derawan, Cm(27.4), Cm(12.9))
        slide2.shapes.add_picture(ConShape_derawan, Cm(29.8), Cm(12.9))
        slide2.shapes.add_picture(EngShape_derawan, Cm(32.1), Cm(12.9))

        # slide3 Bulk Karimun
        slide3.placeholders[10].text = 'EQUIPMENT DASHBOARD \n Bulk Karimun'
        slide3.shapes.add_picture('img/crane_1 Bulk Karimun.png', Cm(0.75), Cm(3.15), Cm(13.75), Cm(4))
        slide3.shapes.add_picture('img/crane_2 Bulk Karimun.png', Cm(0.75), Cm(7.4), Cm(13.75), Cm(4))
        slide3.shapes.add_picture('img/toh Bulk Karimun.png', Cm(14.75), Cm(3.15), Cm(11), Cm(4))
        slide3.shapes.add_picture('img/goh Bulk Karimun.png', Cm(14.75), Cm(7.4), Cm(11), Cm(4))
        slide3.shapes.add_picture('img/cbu Bulk Karimun.png', Cm(0.75), Cm(11.65), Cm(13.75), Cm(5.5))
        slide3.shapes.add_picture('img/au Bulk Karimun.png', Cm(14.75), Cm(11.65), Cm(11), Cm(5.5))
        slide3.shapes.add_picture('img/cs_conveyor Bulk Karimun.png', Cm(28.5), Cm(14.2), Cm(2), Cm(2.5))


        slide3.placeholders[12].text = 'Conveyor-Maker Check-Up'
        slide3.placeholders[13].text = 'Crane-Maker Check-Up'
        slide3.placeholders[14].text = 'Engine-Maker Check-Up'
        slide3.placeholders[15].text = 'Others-Shi ELD'
        slide3.placeholders[16].text = 'Next Docking \n Intermediate \n Survey (IS)'
        slide3.placeholders[17].text = 'Conveyor Belt Spares'
        slide3.placeholders[18].text = 'Crane-Wire Rope Spares'
        slide3.placeholders[19].text = 'Crane'
        slide3.placeholders[20].text = 'Conv.'
        slide3.placeholders[21].text = 'Eng.'

        cal = slide3.shapes.add_picture(NextDockingCol_karimun, Cm(30), Cm(7.7))

        thn = slide3.shapes.add_textbox(Cm(30.6), Cm(8), 1, 0.5)
        thn = thn.text_frame.paragraphs[0]
        thn.text = t_karimun.strftime("%Y")
        font = thn.font
        font.bold = True
        font.name = 'Arial Black'
        font.size = Pt(14)

        tgl = slide3.shapes.add_textbox(Cm(30.8), Cm(8.6), 1, 0.5)
        tgl = tgl.text_frame.paragraphs[0]
        tgl.text = t_karimun.strftime("%d")
        font = tgl.font
        font.bold = True
        font.name = 'Arial Black'
        font.size = Pt(18)

        bln = slide3.shapes.add_textbox(Cm(30.8), Cm(9.35), 1, 0.5)
        bln = bln.text_frame.paragraphs[0]
        bln.text = t_karimun.strftime("%b")
        font = bln.font
        font.name = 'Arial'
        font.size = Pt(14)

        slide3.shapes.add_picture(ConMakerShape_karimun, Cm(32.1), Cm(3.3))
        slide3.shapes.add_picture(CraneMakerShape_karimun, Cm(32.1), Cm(4.4))
        slide3.shapes.add_picture(EngMakerShape_karimun, Cm(32.1), Cm(5.5))
        slide3.shapes.add_picture(OthersShape_karimun, Cm(32.1), Cm(6.6))
        slide3.shapes.add_picture(ConSpareShape_karimun, Cm(32.1), Cm(10.7))
        slide3.shapes.add_picture(CraneSpareShape_karimun, Cm(32.1), Cm(11.8))
        slide3.shapes.add_picture(CraneShape_karimun, Cm(27.4), Cm(12.9))
        slide3.shapes.add_picture(ConShape_karimun, Cm(29.8), Cm(12.9))
        slide3.shapes.add_picture(EngShape_karimun, Cm(32.1), Cm(12.9))

        # slide4 Bulk Dewata
        slide4.placeholders[10].text = 'EQUIPMENT DASHBOARD \n Bulk Dewata'
        slide4.shapes.add_picture('img/hou_c1 Bulk Dewata.png', Cm(0.75), Cm(3.15), Cm(4), Cm(4))
        slide4.shapes.add_picture('img/hou_c2 Bulk Dewata.png', Cm(0.75), Cm(7.4), Cm(4), Cm(4))
        slide4.shapes.add_picture('img/crane_1 Bulk Dewata.png', Cm(5), Cm(3.15), Cm(9.5), Cm(4))
        slide4.shapes.add_picture('img/crane_2 Bulk Dewata.png', Cm(5), Cm(7.4), Cm(9.5), Cm(4))
        slide4.shapes.add_picture('img/toh Bulk Dewata.png', Cm(14.75), Cm(3.15), Cm(11), Cm(4))
        slide4.shapes.add_picture('img/goh Bulk Dewata.png', Cm(14.75), Cm(7.4), Cm(11), Cm(4))
        slide4.shapes.add_picture('img/cbu Bulk Dewata.png', Cm(0.75), Cm(11.65), Cm(13.75), Cm(5.5))
        slide4.shapes.add_picture('img/au Bulk Dewata.png', Cm(14.75), Cm(11.65), Cm(11), Cm(5.5))
        slide4.shapes.add_picture('img/cs_crane Bulk Dewata.png', Cm(26.2), Cm(14.2), Cm(2), Cm(2.5))
        slide4.shapes.add_picture('img/cs_conveyor Bulk Dewata.png', Cm(28.5), Cm(14.2), Cm(2), Cm(2.5))
        slide4.shapes.add_picture('img/cs_engine Bulk Dewata.png', Cm(30.9), Cm(14.2), Cm(2), Cm(2.5))


        slide4.placeholders[12].text = 'Conveyor-Maker Check-Up'
        slide4.placeholders[13].text = 'Crane-Maker Check-Up'
        slide4.placeholders[14].text = 'Engine-Maker Check-Up'
        slide4.placeholders[15].text = 'Others-Shi ELD'
        slide4.placeholders[16].text = 'Next Docking \n Intermediate \n Survey (IS)'
        slide4.placeholders[17].text = 'Conveyor Belt Spares'
        slide4.placeholders[18].text = 'Crane-Wire Rope Spares'
        slide4.placeholders[19].text = 'Crane'
        slide4.placeholders[20].text = 'Conv.'
        slide4.placeholders[21].text = 'Eng.'

        cal = slide4.shapes.add_picture(NextDockingCol_dewata, Cm(30), Cm(7.7))

        thn = slide4.shapes.add_textbox(Cm(30.6), Cm(8), 1, 0.5)
        thn = thn.text_frame.paragraphs[0]
        thn.text = t_dewata.strftime("%Y")
        font = thn.font
        font.bold = True
        font.name = 'Arial Black'
        font.size = Pt(14)

        tgl = slide4.shapes.add_textbox(Cm(30.8), Cm(8.6), 1, 0.5)
        tgl = tgl.text_frame.paragraphs[0]
        tgl.text = t_dewata.strftime("%d")
        font = tgl.font
        font.bold = True
        font.name = 'Arial Black'
        font.size = Pt(18)

        bln = slide4.shapes.add_textbox(Cm(30.8), Cm(9.35), 1, 0.5)
        bln = bln.text_frame.paragraphs[0]
        bln.text = t_dewata.strftime("%b")
        font = bln.font
        font.name = 'Arial'
        font.size = Pt(14)

        slide4.shapes.add_picture(ConMakerShape_dewata, Cm(32.1), Cm(3.3))
        slide4.shapes.add_picture(CraneMakerShape_dewata, Cm(32.1), Cm(4.4))
        slide4.shapes.add_picture(EngMakerShape_dewata, Cm(32.1), Cm(5.5))
        slide4.shapes.add_picture(OthersShape_dewata, Cm(32.1), Cm(6.6))
        slide4.shapes.add_picture(ConSpareShape_dewata, Cm(32.1), Cm(10.7))
        slide4.shapes.add_picture(CraneSpareShape_dewata, Cm(32.1), Cm(11.8))
        slide4.shapes.add_picture(CraneShape_dewata, Cm(27.4), Cm(12.9))
        slide4.shapes.add_picture(ConShape_dewata, Cm(29.8), Cm(12.9))
        slide4.shapes.add_picture(EngShape_dewata, Cm(32.1), Cm(12.9))

        # slide5 Bulk Sumba
        slide5.placeholders[10].text = 'EQUIPMENT DASHBOARD \n Bulk Sumba'
        slide5.shapes.add_picture('img/hou_c1 Bulk Sumba.png', Cm(0.75), Cm(3.15), Cm(4), Cm(6.875))
        slide5.shapes.add_picture('img/hou_c2 Bulk Sumba.png', Cm(0.75), Cm(10.275), Cm(4), Cm(6.875))
        slide5.shapes.add_picture('img/crane_1 Bulk Sumba.png', Cm(5), Cm(3.15), Cm(9.5), Cm(6.875))
        slide5.shapes.add_picture('img/crane_2 Bulk Sumba.png', Cm(5), Cm(10.275), Cm(9.5), Cm(6.875))
        slide5.shapes.add_picture('img/toh Bulk Sumba.png', Cm(14.75), Cm(3.15), Cm(11), Cm(4))
        slide5.shapes.add_picture('img/goh Bulk Sumba.png', Cm(14.75), Cm(7.4), Cm(11), Cm(4))
        slide5.shapes.add_picture('img/au Bulk Sumba.png', Cm(14.75), Cm(11.65), Cm(11), Cm(5.5))
        slide5.shapes.add_picture('img/cs_conveyor Bulk Sumba.png', Cm(28.5), Cm(14.2), Cm(2), Cm(2.5))
        slide5.shapes.add_picture('img/cs_engine Bulk Sumba.png', Cm(30.9), Cm(14.2), Cm(2), Cm(2.5))


        slide5.placeholders[12].text = 'Conveyor-Maker Check-Up'
        slide5.placeholders[13].text = 'Crane-Maker Check-Up'
        slide5.placeholders[14].text = 'Engine-Maker Check-Up'
        slide5.placeholders[15].text = 'Others-Shi ELD'
        slide5.placeholders[16].text = 'Next Docking \n Intermediate \n Survey (IS)'
        slide5.placeholders[17].text = 'Conveyor Belt Spares'
        slide5.placeholders[18].text = 'Crane-Wire Rope Spares'
        slide5.placeholders[19].text = 'Crane'
        slide5.placeholders[20].text = 'Conv.'
        slide5.placeholders[21].text = 'Eng.'

        cal = slide5.shapes.add_picture(NextDockingCol_sumba, Cm(30), Cm(7.7))

        thn = slide5.shapes.add_textbox(Cm(30.6), Cm(8), 1, 0.5)
        thn = thn.text_frame.paragraphs[0]
        thn.text = t_sumba.strftime("%Y")
        font = thn.font
        font.bold = True
        font.name = 'Arial Black'
        font.size = Pt(14)

        tgl = slide5.shapes.add_textbox(Cm(30.8), Cm(8.6), 1, 0.5)
        tgl = tgl.text_frame.paragraphs[0]
        tgl.text = t_sumba.strftime("%d")
        font = tgl.font
        font.bold = True
        font.name = 'Arial Black'
        font.size = Pt(18)

        bln = slide5.shapes.add_textbox(Cm(30.8), Cm(9.35), 1, 0.5)
        bln = bln.text_frame.paragraphs[0]
        bln.text = t_sumba.strftime("%b")
        font = bln.font
        font.name = 'Arial'
        font.size = Pt(14)

        slide5.shapes.add_picture(ConMakerShape_sumba, Cm(32.1), Cm(3.3))
        slide5.shapes.add_picture(CraneMakerShape_sumba, Cm(32.1), Cm(4.4))
        slide5.shapes.add_picture(EngMakerShape_sumba, Cm(32.1), Cm(5.5))
        slide5.shapes.add_picture(OthersShape_sumba, Cm(32.1), Cm(6.6))
        slide5.shapes.add_picture(ConSpareShape_sumba, Cm(32.1), Cm(10.7))
        slide5.shapes.add_picture(CraneSpareShape_sumba, Cm(32.1), Cm(11.8))
        slide5.shapes.add_picture(CraneShape_sumba, Cm(27.4), Cm(12.9))
        slide5.shapes.add_picture(ConShape_sumba, Cm(29.8), Cm(12.9))
        slide5.shapes.add_picture(EngShape_sumba, Cm(32.1), Cm(12.9))
        
        # slide6 Bulk Java
        slide6.placeholders[10].text = 'EQUIPMENT DASHBOARD \n Bulk Java'
        slide6.shapes.add_picture('img/hou_c1 Bulk Java.png', Cm(0.75), Cm(3.15), Cm(4), Cm(4))
        slide6.shapes.add_picture('img/hou_c2 Bulk Java.png', Cm(0.75), Cm(7.4), Cm(4), Cm(4))
        slide6.shapes.add_picture('img/crane_1 Bulk Java.png', Cm(5), Cm(3.15), Cm(9.5), Cm(4))
        slide6.shapes.add_picture('img/crane_2 Bulk Java.png', Cm(5), Cm(7.4), Cm(9.5), Cm(4))
        slide6.shapes.add_picture('img/toh Bulk Java.png', Cm(14.75), Cm(3.15), Cm(11), Cm(4))
        slide6.shapes.add_picture('img/goh Bulk Java.png', Cm(14.75), Cm(7.4), Cm(11), Cm(4))
        slide6.shapes.add_picture('img/cbu Bulk Java.png', Cm(0.75), Cm(11.65), Cm(13.75), Cm(5.5))
        slide6.shapes.add_picture('img/au Bulk Java.png', Cm(14.75), Cm(11.65), Cm(11), Cm(5.5))
        slide6.shapes.add_picture('img/cs_crane Bulk Java.png', Cm(26.2), Cm(14.2), Cm(2), Cm(2.5))
        slide6.shapes.add_picture('img/cs_conveyor Bulk Java.png', Cm(28.5), Cm(14.2), Cm(2), Cm(2.5))
        slide6.shapes.add_picture('img/cs_engine Bulk Java.png', Cm(30.9), Cm(14.2), Cm(2), Cm(2.5))


        slide6.placeholders[12].text = 'Conveyor-Maker Check-Up'
        slide6.placeholders[13].text = 'Crane-Maker Check-Up'
        slide6.placeholders[14].text = 'Engine-Maker Check-Up'
        slide6.placeholders[15].text = 'Others-Shi ELD'
        slide6.placeholders[16].text = 'Next Docking \n Intermediate \n Survey (IS)'
        slide6.placeholders[17].text = 'Conveyor Belt Spares'
        slide6.placeholders[18].text = 'Crane-Wire Rope Spares'
        slide6.placeholders[19].text = 'Crane'
        slide6.placeholders[20].text = 'Conv.'
        slide6.placeholders[21].text = 'Eng.'

        cal = slide6.shapes.add_picture(NextDockingCol_java, Cm(30), Cm(7.7))

        thn = slide6.shapes.add_textbox(Cm(30.6), Cm(8), 1, 0.5)
        thn = thn.text_frame.paragraphs[0]
        thn.text = t_java.strftime("%Y")
        font = thn.font
        font.bold = True
        font.name = 'Arial Black'
        font.size = Pt(14)

        tgl = slide6.shapes.add_textbox(Cm(30.8), Cm(8.6), 1, 0.5)
        tgl = tgl.text_frame.paragraphs[0]
        tgl.text = t_java.strftime("%d")
        font = tgl.font
        font.bold = True
        font.name = 'Arial Black'
        font.size = Pt(18)

        bln = slide6.shapes.add_textbox(Cm(30.8), Cm(9.35), 1, 0.5)
        bln = bln.text_frame.paragraphs[0]
        bln.text = t_java.strftime("%b")
        font = bln.font
        font.name = 'Arial'
        font.size = Pt(14)

        slide6.shapes.add_picture(ConMakerShape_java, Cm(32.1), Cm(3.3))
        slide6.shapes.add_picture(CraneMakerShape_java, Cm(32.1), Cm(4.4))
        slide6.shapes.add_picture(EngMakerShape_java, Cm(32.1), Cm(5.5))
        slide6.shapes.add_picture(OthersShape_java, Cm(32.1), Cm(6.6))
        slide6.shapes.add_picture(ConSpareShape_java, Cm(32.1), Cm(10.7))
        slide6.shapes.add_picture(CraneSpareShape_java, Cm(32.1), Cm(11.8))
        slide6.shapes.add_picture(CraneShape_java, Cm(27.4), Cm(12.9))
        slide6.shapes.add_picture(ConShape_java, Cm(29.8), Cm(12.9))
        slide6.shapes.add_picture(EngShape_java, Cm(32.1), Cm(12.9))

        # slide7 Bulk Java
        slide7.placeholders[10].text = 'EQUIPMENT DASHBOARD \n Bulk Borneo'
        slide7.shapes.add_picture('img/hou_c1 Bulk Borneo.png', Cm(0.75), Cm(3.15), Cm(4), Cm(4))
        slide7.shapes.add_picture('img/hou_c2 Bulk Borneo.png', Cm(0.75), Cm(7.4), Cm(4), Cm(4))
        slide7.shapes.add_picture('img/crane_1 Bulk Borneo.png', Cm(5), Cm(3.15), Cm(9.5), Cm(4))
        slide7.shapes.add_picture('img/crane_2 Bulk Borneo.png', Cm(5), Cm(7.4), Cm(9.5), Cm(4))
        slide7.shapes.add_picture('img/toh Bulk Borneo.png', Cm(14.75), Cm(3.15), Cm(11), Cm(4))
        slide7.shapes.add_picture('img/goh Bulk Borneo.png', Cm(14.75), Cm(7.4), Cm(11), Cm(4))
        slide7.shapes.add_picture('img/cbu Bulk Borneo.png', Cm(0.75), Cm(11.65), Cm(13.75), Cm(5.5))
        slide7.shapes.add_picture('img/au Bulk Borneo.png', Cm(14.75), Cm(11.65), Cm(11), Cm(5.5))
        slide7.shapes.add_picture('img/cs_crane Bulk Borneo.png', Cm(26.2), Cm(14.2), Cm(2), Cm(2.5))
        slide7.shapes.add_picture('img/cs_conveyor Bulk Borneo.png', Cm(28.5), Cm(14.2), Cm(2), Cm(2.5))
        slide7.shapes.add_picture('img/cs_engine Bulk Borneo.png', Cm(30.9), Cm(14.2), Cm(2), Cm(2.5))


        slide7.placeholders[12].text = 'Conveyor-Maker Check-Up'
        slide7.placeholders[13].text = 'Crane-Maker Check-Up'
        slide7.placeholders[14].text = 'Engine-Maker Check-Up'
        slide7.placeholders[15].text = 'Others-Shi ELD'
        slide7.placeholders[16].text = 'Next Docking \n Intermediate \n Survey (IS)'
        slide7.placeholders[17].text = 'Conveyor Belt Spares'
        slide7.placeholders[18].text = 'Crane-Wire Rope Spares'
        slide7.placeholders[19].text = 'Crane'
        slide7.placeholders[20].text = 'Conv.'
        slide7.placeholders[21].text = 'Eng.'

        cal = slide7.shapes.add_picture(NextDockingCol_borneo, Cm(30), Cm(7.7))

        thn = slide7.shapes.add_textbox(Cm(30.6), Cm(8), 1, 0.5)
        thn = thn.text_frame.paragraphs[0]
        thn.text = t_borneo.strftime("%Y")
        font = thn.font
        font.bold = True
        font.name = 'Arial Black'
        font.size = Pt(14)

        tgl = slide7.shapes.add_textbox(Cm(30.8), Cm(8.6), 1, 0.5)
        tgl = tgl.text_frame.paragraphs[0]
        tgl.text = t_borneo.strftime("%d")
        font = tgl.font
        font.bold = True
        font.name = 'Arial Black'
        font.size = Pt(18)

        bln = slide7.shapes.add_textbox(Cm(30.8), Cm(9.35), 1, 0.5)
        bln = bln.text_frame.paragraphs[0]
        bln.text = t_borneo.strftime("%b")
        font = bln.font
        font.name = 'Arial'
        font.size = Pt(14)

        slide7.shapes.add_picture(ConMakerShape_borneo, Cm(32.1), Cm(3.3))
        slide7.shapes.add_picture(CraneMakerShape_borneo, Cm(32.1), Cm(4.4))
        slide7.shapes.add_picture(EngMakerShape_borneo, Cm(32.1), Cm(5.5))
        slide7.shapes.add_picture(OthersShape_borneo, Cm(32.1), Cm(6.6))
        slide7.shapes.add_picture(ConSpareShape_borneo, Cm(32.1), Cm(10.7))
        slide7.shapes.add_picture(CraneSpareShape_borneo, Cm(32.1), Cm(11.8))
        slide7.shapes.add_picture(CraneShape_borneo, Cm(27.4), Cm(12.9))
        slide7.shapes.add_picture(ConShape_borneo, Cm(29.8), Cm(12.9))
        slide7.shapes.add_picture(EngShape_borneo, Cm(32.1), Cm(12.9))

        # slide8 Bulk Java
        slide8.placeholders[10].text = 'EQUIPMENT DASHBOARD \n Bulk Celebes'
        slide8.shapes.add_picture('img/hou_c1 Bulk Celebes.png', Cm(0.75), Cm(3.15), Cm(4), Cm(4))
        slide8.shapes.add_picture('img/hou_c2 Bulk Celebes.png', Cm(0.75), Cm(7.4), Cm(4), Cm(4))
        slide8.shapes.add_picture('img/crane_1 Bulk Celebes.png', Cm(5), Cm(3.15), Cm(9.5), Cm(4))
        slide8.shapes.add_picture('img/crane_2 Bulk Celebes.png', Cm(5), Cm(7.4), Cm(9.5), Cm(4))
        slide8.shapes.add_picture('img/toh Bulk Celebes.png', Cm(14.75), Cm(3.15), Cm(11), Cm(4))
        slide8.shapes.add_picture('img/goh Bulk Celebes.png', Cm(14.75), Cm(7.4), Cm(11), Cm(4))
        slide8.shapes.add_picture('img/cbu Bulk Celebes.png', Cm(0.75), Cm(11.65), Cm(13.75), Cm(5.5))
        slide8.shapes.add_picture('img/au Bulk Celebes.png', Cm(14.75), Cm(11.65), Cm(11), Cm(5.5))
        slide8.shapes.add_picture('img/cs_crane Bulk Celebes.png', Cm(26.2), Cm(14.2), Cm(2), Cm(2.5))
        slide8.shapes.add_picture('img/cs_conveyor Bulk Celebes.png', Cm(28.5), Cm(14.2), Cm(2), Cm(2.5))
        slide8.shapes.add_picture('img/cs_engine Bulk Celebes.png', Cm(30.9), Cm(14.2), Cm(2), Cm(2.5))


        slide8.placeholders[12].text = 'Conveyor-Maker Check-Up'
        slide8.placeholders[13].text = 'Crane-Maker Check-Up'
        slide8.placeholders[14].text = 'Engine-Maker Check-Up'
        slide8.placeholders[15].text = 'Others-Shi ELD'
        slide8.placeholders[16].text = 'Next Docking \n Intermediate \n Survey (IS)'
        slide8.placeholders[17].text = 'Conveyor Belt Spares'
        slide8.placeholders[18].text = 'Crane-Wire Rope Spares'
        slide8.placeholders[19].text = 'Crane'
        slide8.placeholders[20].text = 'Conv.'
        slide8.placeholders[21].text = 'Eng.'

        cal = slide8.shapes.add_picture(NextDockingCol_celebes, Cm(30), Cm(7.7))

        thn = slide8.shapes.add_textbox(Cm(30.6), Cm(8), 1, 0.5)
        thn = thn.text_frame.paragraphs[0]
        thn.text = t_celebes.strftime("%Y")
        font = thn.font
        font.bold = True
        font.name = 'Arial Black'
        font.size = Pt(14)

        tgl = slide8.shapes.add_textbox(Cm(30.8), Cm(8.6), 1, 0.5)
        tgl = tgl.text_frame.paragraphs[0]
        tgl.text = t_celebes.strftime("%d")
        font = tgl.font
        font.bold = True
        font.name = 'Arial Black'
        font.size = Pt(18)

        bln = slide8.shapes.add_textbox(Cm(30.8), Cm(9.35), 1, 0.5)
        bln = bln.text_frame.paragraphs[0]
        bln.text = t_celebes.strftime("%b")
        font = bln.font
        font.name = 'Arial'
        font.size = Pt(14)

        slide8.shapes.add_picture(ConMakerShape_celebes, Cm(32.1), Cm(3.3))
        slide8.shapes.add_picture(CraneMakerShape_celebes, Cm(32.1), Cm(4.4))
        slide8.shapes.add_picture(EngMakerShape_celebes, Cm(32.1), Cm(5.5))
        slide8.shapes.add_picture(OthersShape_celebes, Cm(32.1), Cm(6.6))
        slide8.shapes.add_picture(ConSpareShape_celebes, Cm(32.1), Cm(10.7))
        slide8.shapes.add_picture(CraneSpareShape_celebes, Cm(32.1), Cm(11.8))
        slide8.shapes.add_picture(CraneShape_celebes, Cm(27.4), Cm(12.9))
        slide8.shapes.add_picture(ConShape_celebes, Cm(29.8), Cm(12.9))
        slide8.shapes.add_picture(EngShape_celebes, Cm(32.1), Cm(12.9))

        # slide9 Bulk Java
        slide9.placeholders[10].text = 'EQUIPMENT DASHBOARD \n Bulk Natuna'
        slide9.shapes.add_picture('img/hou_c1 Bulk Natuna.png', Cm(0.75), Cm(3.15), Cm(4), Cm(14))
        slide9.shapes.add_picture('img/crane_1 Bulk Natuna.png', Cm(5), Cm(3.15), Cm(9.5), Cm(14))
        slide9.shapes.add_picture('img/toh Bulk Natuna.png', Cm(14.75), Cm(3.15), Cm(11), Cm(4))
        slide9.shapes.add_picture('img/goh Bulk Natuna.png', Cm(14.75), Cm(7.4), Cm(11), Cm(4))
        slide9.shapes.add_picture('img/au Bulk Natuna.png', Cm(14.75), Cm(11.65), Cm(11), Cm(5.5))


        slide9.placeholders[12].text = 'Conveyor-Maker Check-Up'
        slide9.placeholders[13].text = 'Crane-Maker Check-Up'
        slide9.placeholders[14].text = 'Engine-Maker Check-Up'
        slide9.placeholders[15].text = 'Others-Shi ELD'
        slide9.placeholders[16].text = 'Next Docking \n Intermediate \n Survey (IS)'
        slide9.placeholders[17].text = 'Conveyor Belt Spares'
        slide9.placeholders[18].text = 'Crane-Wire Rope Spares'
        slide9.placeholders[19].text = 'Crane'
        slide9.placeholders[20].text = 'Conv.'
        slide9.placeholders[21].text = 'Eng.'

        cal = slide9.shapes.add_picture(NextDockingCol_natuna, Cm(30), Cm(7.7))

        thn = slide9.shapes.add_textbox(Cm(30.6), Cm(8), 1, 0.5)
        thn = thn.text_frame.paragraphs[0]
        thn.text = t_natuna.strftime("%Y")
        font = thn.font
        font.bold = True
        font.name = 'Arial Black'
        font.size = Pt(14)

        tgl = slide9.shapes.add_textbox(Cm(30.8), Cm(8.6), 1, 0.5)
        tgl = tgl.text_frame.paragraphs[0]
        tgl.text = t_natuna.strftime("%d")
        font = tgl.font
        font.bold = True
        font.name = 'Arial Black'
        font.size = Pt(18)

        bln = slide9.shapes.add_textbox(Cm(30.8), Cm(9.35), 1, 0.5)
        bln = bln.text_frame.paragraphs[0]
        bln.text = t_natuna.strftime("%b")
        font = bln.font
        font.name = 'Arial'
        font.size = Pt(14)

        slide9.shapes.add_picture(ConMakerShape_natuna, Cm(32.1), Cm(3.3))
        slide9.shapes.add_picture(CraneMakerShape_natuna, Cm(32.1), Cm(4.4))
        slide9.shapes.add_picture(EngMakerShape_natuna, Cm(32.1), Cm(5.5))
        slide9.shapes.add_picture(OthersShape_natuna, Cm(32.1), Cm(6.6))
        slide9.shapes.add_picture(ConSpareShape_natuna, Cm(32.1), Cm(10.7))
        slide9.shapes.add_picture(CraneSpareShape_natuna, Cm(32.1), Cm(11.8))
        slide9.shapes.add_picture(CraneShape_natuna, Cm(27.4), Cm(12.9))
        slide9.shapes.add_picture(ConShape_natuna, Cm(29.8), Cm(12.9))
        slide9.shapes.add_picture(EngShape_natuna, Cm(32.1), Cm(12.9))

        prs.save(bytes_io)    

        # return dcc.send_bytes(to_pptx, 'Equipment Dashboard.pptx')
    
    return dcc.send_bytes(to_pptx, 'Equipment Dashboard.pptx'), download_date


######################
# Plot Function
######################
##-----Function
#-- 1. Plot Hydrolic Utilization
def plot_hou1(value):
    color=''
    if value<=(100/3) : color='#488f31'
    elif value<=(100/3*2) : color='#ffe600'
    elif value<=100 : color='#de425b'

    fig = go.Figure(go.Indicator(
        mode = "gauge+number",
        value = value,
        number = {'suffix': "%" },
        domain = {'x': [0.1, 0.9], 'y': [0.03, 0.5]},
        gauge = {
            'axis': {'range': [None, 100], 'tickwidth': 0, 'tickcolor': "white", 'visible': False},
            'bgcolor': "white",
            'borderwidth': 3,
            'bordercolor': color,
            'bar':{'color':color},
            }))
    fig.add_annotation(x=0.5, y=0.85, text='<b>Hydrolic Oil Utilization</b><br>Crane no 1<br>(Max 6,000hrs)', 
                       font=dict(size=57), showarrow=False)
    fig.update_layout({'height':700,'width':750,'plot_bgcolor': '#ffffff', 'paper_bgcolor': '#ffffff',
                       'margin' : {'t':3, 'b':3, 'l':3, 'r':3}, "autosize": True})
    return(fig)

def plot_hou2(value):
    color=''
    if value<=(100/3) : color='#488f31'
    elif value<=(100/3*2) : color='#ffe600'
    elif value<=100 : color='#de425b'

    fig = go.Figure(go.Indicator(
        mode = "gauge+number",
        value = value,
        number = {'suffix': "%" },
        domain = {'x': [0.1, 0.9], 'y': [0.03, 0.5]},
        gauge = {
            'axis': {'range': [None, 100], 'tickwidth': 0, 'tickcolor': "white", 'visible': False},
            'bgcolor': "white",
            'borderwidth': 3,
            'bordercolor': color,
            'bar':{'color':color},
            }))
    fig.add_annotation(x=0.5, y=0.85, text='<b>Hydrolic Oil Utilization</b><br>Crane no 2<br>(Max 6,000hrs)', 
                       font=dict(size=57), showarrow=False)
    fig.update_layout({'height':700,'width':750,'plot_bgcolor': '#ffffff', 'paper_bgcolor': '#ffffff',
                       'margin' : {'t':3, 'b':3, 'l':3, 'r':3}, "autosize": True})
    return(fig)


#-- 1. Plot Hydrolic Utilization (Bulk Sumba)
def plot_hou1_s(value):
    color=''
    if value<=(100/3) : color='#488f31'
    elif value<=(100/3*2) : color='#ffe600'
    elif value<=100 : color='#de425b'

    fig = go.Figure(go.Indicator(
        mode = "gauge+number",
        value = value,
        number = {'suffix': "%" },
        domain = {'x': [0.1, 0.9], 'y': [0.03, 0.5]},
        gauge = {
            'axis': {'range': [None, 100], 'tickwidth': 0, 'tickcolor': "white", 'visible': False},
            'bgcolor': "white",
            'borderwidth': 3,
            'bordercolor': color,
            'bar':{'color':color},
            }))
    fig.add_annotation(x=0.5, y=0.7, text='<b>Hydrolic Oil Utilization</b><br>Crane no 1<br>(Max 6,000hrs)', 
                       font=dict(size=57), showarrow=False)
    fig.update_layout({'height':1200,'width':750,'plot_bgcolor': '#ffffff', 'paper_bgcolor': '#ffffff',
                       'margin' : {'t':3, 'b':150, 'l':3, 'r':3}, "autosize": True})
    return(fig)

def plot_hou2_s(value):
    color=''
    if value<=(100/3) : color='#488f31'
    elif value<=(100/3*2) : color='#ffe600'
    elif value<=100 : color='#de425b'

    fig = go.Figure(go.Indicator(
        mode = "gauge+number",
        value = value,
        number = {'suffix': "%" },
        domain = {'x': [0.1, 0.9], 'y': [0.03, 0.5]},
        gauge = {
            'axis': {'range': [None, 100], 'tickwidth': 0, 'tickcolor': "white", 'visible': False},
            'bgcolor': "white",
            'borderwidth': 3,
            'bordercolor': color,
            'bar':{'color':color},
            }))
    fig.add_annotation(x=0.5, y=0.7, text='<b>Hydrolic Oil Utilization</b><br>Crane no 2<br>(Max 6,000hrs)', 
                       font=dict(size=57), showarrow=False)
    fig.update_layout({'height':1200,'width':750,'plot_bgcolor': '#ffffff', 'paper_bgcolor': '#ffffff',
                       'margin' : {'t':3, 'b':150, 'l':3, 'r':3}, "autosize": True})
    return(fig)


#-- 1. Plot Hydrolic Utilization
def plot_hou1_n(value):
    color=''
    if value<=(100/3) : color='#488f31'
    elif value<=(100/3*2) : color='#ffe600'
    elif value<=100 : color='#de425b'

    fig = go.Figure(go.Indicator(
        mode = "gauge+number",
        value = value,
        number = {'suffix': "%" },
        domain = {'x': [0.1, 0.9], 'y': [0.03, 0.5]},
        gauge = {
            'axis': {'range': [None, 100], 'tickwidth': 0, 'tickcolor': "white", 'visible': False},
            'bgcolor': "white",
            'borderwidth': 3,
            'bordercolor': color,
            'bar':{'color':color},
            }))
    fig.add_annotation(x=0.5, y=0.5, text='<b>Hydrolic Oil Utilization</b><br>Crane no 1<br>(Max 6,000hrs)', 
                       font=dict(size=57), showarrow=False)
    fig.update_layout({'height':2500,'width':750,'plot_bgcolor': '#ffffff', 'paper_bgcolor': '#ffffff',
                       'margin' : {'t':3, 'b':300, 'l':3, 'r':3}, "autosize": True})
    return(fig)

#-- 2. Plot TOH & GOH
def plot_toh(value):
    color=[]
    for i in value:
        if i<=20 : color.append('#488f31')
        elif i<=40 : color.append('#88c580')
        elif i<=60 : color.append('#ffe600')
        elif i<=80 : color.append('#f7a258')
        else: color.append('#de425b')
    fig = go.Figure()
    fig.add_trace(go.Indicator(
        mode = "gauge+number",
        value = value[0],
        number = {'suffix': "%" },
        domain = {'x': [0.02, 0.22], 'y': [0, 0.5]},
        title = {'text': "DG01", 'font': {'size': 60}},
        gauge = {
            'axis': {'range': [None, 100], 'tickwidth': 0, 'tickcolor': "white", 'visible': False},
            'bgcolor': "white",
            'borderwidth': 3,
            'bordercolor': color[0],
            'bar':{'color':color[0]},
            }))
    fig.add_trace(go.Indicator(
        mode = "gauge+number",
        value = value[1],
        number = {'suffix': "%" },
        domain = {'x': [0.27, 0.47], 'y': [0, 0.5]},
        title = {'text': "DG02", 'font': {'size': 60}},
        gauge = {
            'axis': {'range': [None, 100], 'tickwidth': 0, 'tickcolor': "white", 'visible': False},
            'bgcolor': "white",
            'borderwidth': 3,
            'bordercolor': color[1],
            'bar':{'color':color[1]},
            }))
    fig.add_trace(go.Indicator(
        mode = "gauge+number",
        value = value[2],
        number = {'suffix': "%" },
        domain = {'x': [0.52, 0.72], 'y': [0, 0.5]},
        title = {'text': "DG03", 'font': {'size': 60}},
        gauge = {
            'axis': {'range': [None, 100], 'tickwidth': 0, 'tickcolor': "white", 'visible': False},
            'bgcolor': "white",
            'borderwidth': 3,
            'bordercolor': color[2],
            'bar':{'color':color[2]},
            }))
    fig.add_trace(go.Indicator(
        mode = "gauge+number",
        value = value[3],
        number = {'suffix': "%" },
        domain = {'x': [0.77, 0.97], 'y': [0, 0.5]},
        title = {'text': "DG04", 'font': {'size': 60}},
        gauge = {
            'axis': {'range': [None, 100], 'tickwidth': 0, 'tickcolor': "white", 'visible': False},
            'bgcolor': "white",
            'borderwidth': 3,
            'bordercolor': color[3],
            'bar':{'color':color[3]},
            }))
    fig.add_annotation(x=0.5, y=0.90, text='<b>Engine Overhaul</b><br>TOH to be conducted @12,000hrs<br>', 
                       font=dict(size=75), showarrow=False)
    fig.update_layout({'height':700,'width':1900,'plot_bgcolor': '#ffffff', 'paper_bgcolor': '#ffffff',
                       'margin' : {'t':3, 'b':3, 'l':3, 'r':3}, "autosize": True})
    

    return(fig)

############################################################################################################################

def plot_goh(value):
    color=[]
    for i in value:
        if i<=20 : color.append('#488f31')
        elif i<=40 : color.append('#88c580')
        elif i<=60 : color.append('#ffe600')
        elif i<=80 : color.append('#f7a258')
        else: color.append('#de425b')
    fig = go.Figure()
    fig.add_trace(go.Indicator(
        mode = "gauge+number",
        value = value[0],
        number = {'suffix': "%" },
        domain = {'x': [0.02, 0.22], 'y': [0, 0.5]},
        title = {'text': "DG01", 'font': {'size': 60}},
        gauge = {
            'axis': {'range': [None, 100], 'tickwidth': 0, 'tickcolor': "white", 'visible': False},
            'bgcolor': "white",
            'borderwidth': 3,
            'bordercolor': color[0],
            'bar':{'color':color[0]},
            }))
    fig.add_trace(go.Indicator(
        mode = "gauge+number",
        value = value[1],
        number = {'suffix': "%" },
        domain = {'x': [0.27, 0.47], 'y': [0, 0.5]},
        title = {'text': "DG02", 'font': {'size': 60}},
        gauge = {
            'axis': {'range': [None, 100], 'tickwidth': 0, 'tickcolor': "white", 'visible': False},
            'bgcolor': "white",
            'borderwidth': 3,
            'bordercolor': color[1],
            'bar':{'color':color[1]},
            }))
    fig.add_trace(go.Indicator(
        mode = "gauge+number",
        value = value[2],
        number = {'suffix': "%" },
        domain = {'x': [0.52, 0.72], 'y': [0, 0.5]},
        title = {'text': "DG03", 'font': {'size': 60}},
        gauge = {
            'axis': {'range': [None, 100], 'tickwidth': 0, 'tickcolor': "white", 'visible': False},
            'bgcolor': "white",
            'borderwidth': 3,
            'bordercolor': color[2],
            'bar':{'color':color[2]},
            }))
    fig.add_trace(go.Indicator(
        mode = "gauge+number",
        value = value[3],
        number = {'suffix': "%" },
        domain = {'x': [0.77, 0.97], 'y': [0, 0.5]},
        title = {'text': "DG04", 'font': {'size': 60}},
        gauge = {
            'axis': {'range': [None, 100], 'tickwidth': 0, 'tickcolor': "white", 'visible': False},
            'bgcolor': "white",
            'borderwidth': 3,
            'bordercolor': color[3],
            'bar':{'color':color[3]},
            }))
    fig.add_annotation(x=0.5, y=0.90, text='<b>Engine Overhaul</b><br>GOH to be conducted @24,000hrs<br>', 
                       font=dict(size=75), showarrow=False)
    fig.update_layout({'height':700,'width':1900,'plot_bgcolor': '#ffffff', 'paper_bgcolor': '#ffffff',
                       'margin' : {'t':3, 'b':3, 'l':3, 'r':3}, "autosize": True})
    

    return(fig)

#-- 3. Plot Wire Rope
def plot_wr1(value, wire_rope_c1):
    anchos=0.5
    color=[]
    position=[]
    for i in value:
        if i<=20 : color.append('#488f31')
        elif i<=40 : color.append('#88c580')
        elif i<=60 : color.append('#ffe600')
        elif i<=80 : color.append('#f7a258')
        else: color.append('#de425b')
    
    for i in value:
        if i<=50 : position.append('outside')
        else : position.append('inside')

    fig = go.Figure()
    fig.add_trace(go.Bar(x = wire_rope_c1['Crane 1'],
                         y = wire_rope_c1['Max'],
                         marker_color='#ededed',
                         hovertemplate='%{y:y}%',
                         width = anchos, name = 'Crane 1'))
    fig.add_trace(go.Bar(x = wire_rope_c1['Crane 1'], 
                         y = value,
                         text=value.apply(lambda x: '{0:1.0f}%'.format(x)),
                         textposition=position,
                         textfont_size=60,
                         textangle=0,
                         hovertemplate='%{y:y}%',
                         marker_color =color,
                         width = anchos, name = 'Crane 1'))
    fig.update_yaxes(visible=False)
    fig.update_layout(title = "<b>Crane 1</b>",
                      title_x=0.5,
                      barmode = 'overlay',
                      title_font_size = 78,
                      showlegend=False,
                      xaxis = dict(tickfont = dict(size=60)))
    fig.update_layout({'height':700,'width':1700,'plot_bgcolor': '#ffffff', 'paper_bgcolor': '#ffffff',
                       'margin' : {'t':250, 'b':3, 'l':3, 'r':3}, "autosize": True
                      })

    return(fig)

############################################################################################################################

def plot_wr2(value, wire_rope_c2):
    anchos = 0.5
    color=[]
    position=[]
    for i in value:
        if i<=20 : color.append('#488f31')
        elif i<=40 : color.append('#88c580')
        elif i<=60 : color.append('#ffe600')
        elif i<=80 : color.append('#f7a258')
        else: color.append('#de425b')
    
    for i in value:
        if i<=50 : position.append('outside')
        else : position.append('inside')

    fig = go.Figure()
    fig.add_trace(go.Bar(x = wire_rope_c2['Crane 2'],
                         y = wire_rope_c2['Max'],
                         marker_color='#ededed',
                         hovertemplate='%{y:y}%',
                         width = anchos, name = 'Crane 2'))
    fig.add_trace(go.Bar(x = wire_rope_c2['Crane 2'], 
                         y = value,
                         text=value.apply(lambda x: '{0:1.0f}%'.format(x)),
                         textposition=position,
                         textfont_size=60,
                         textangle=0,
                         hovertemplate='%{y:y}%',
                         marker_color =color,
                         width = anchos, name = 'Crane 2'))
    fig.update_yaxes(visible=False)
    fig.update_layout(title = "<b>Crane 2</b>",
                      title_x=0.5,
                      barmode = 'overlay',
                      title_font_size = 78,
                      showlegend=False,
                      xaxis = dict(tickfont = dict(size=60)))
    fig.update_layout({'height':700,'width':1700,'plot_bgcolor': '#ffffff', 'paper_bgcolor': '#ffffff',
                       'margin' : {'t':250, 'b':3, 'l':3, 'r':3}, "autosize": True
                      })

    return(fig)


#-- 3. Plot Wire Rope (Bulk Derawan & Bulk Karimun)
def plot_wr1_dk(value, wire_rope_c1):
    anchos=0.5
    color=[]
    position=[]
    for i in value:
        if i<=20 : color.append('#488f31')
        elif i<=40 : color.append('#88c580')
        elif i<=60 : color.append('#ffe600')
        elif i<=80 : color.append('#f7a258')
        else: color.append('#de425b')
    
    for i in value:
        if i<=50 : position.append('outside')
        else : position.append('inside')

    fig = go.Figure()
    fig.add_trace(go.Bar(x = wire_rope_c1['Crane 1'],
                         y = wire_rope_c1['Max'],
                         marker_color='#ededed',
                         hovertemplate='%{y:y}%',
                         width = anchos, name = 'Crane 1'))
    fig.add_trace(go.Bar(x = wire_rope_c1['Crane 1'], 
                         y = value,
                         text=value.apply(lambda x: '{0:1.0f}%'.format(x)),
                         textposition=position,
                         textfont_size=60,
                         textangle=0,
                         hovertemplate='%{y:y}%',
                         marker_color =color,
                         width = anchos, name = 'Crane 1'))
    fig.update_yaxes(visible=False)
    fig.update_layout(title = "<b>Crane 1</b>",
                      title_x=0.5,
                      barmode = 'overlay',
                      title_font_size = 78,
                      showlegend=False,
                      xaxis = dict(tickfont = dict(size=60)))
    fig.update_layout({'height':700,'width':2700,'plot_bgcolor': '#ffffff', 'paper_bgcolor': '#ffffff',
                       'margin' : {'t':250, 'b':3, 'l':3, 'r':3}, "autosize": True
                      })

    return(fig)

############################################################################################################################

def plot_wr2_dk(value, wire_rope_c2):
    anchos = 0.5
    color=[]
    position=[]
    for i in value:
        if i<=20 : color.append('#488f31')
        elif i<=40 : color.append('#88c580')
        elif i<=60 : color.append('#ffe600')
        elif i<=80 : color.append('#f7a258')
        else: color.append('#de425b')
    
    for i in value:
        if i<=50 : position.append('outside')
        else : position.append('inside')

    fig = go.Figure()
    fig.add_trace(go.Bar(x = wire_rope_c2['Crane 2'],
                         y = wire_rope_c2['Max'],
                         marker_color='#ededed',
                         hovertemplate='%{y:y}%',
                         width = anchos, name = 'Crane 2'))
    fig.add_trace(go.Bar(x = wire_rope_c2['Crane 2'], 
                         y = value,
                         text=value.apply(lambda x: '{0:1.0f}%'.format(x)),
                         textposition=position,
                         textfont_size=60,
                         textangle=0,
                         hovertemplate='%{y:y}%',
                         marker_color =color,
                         width = anchos, name = 'Crane 2'))
    fig.update_yaxes(visible=False)
    fig.update_layout(title = "<b>Crane 2</b>",
                      title_x=0.5,
                      barmode = 'overlay',
                      title_font_size = 78,
                      showlegend=False,
                      xaxis = dict(tickfont = dict(size=60)))
    fig.update_layout({'height':700,'width':2700,'plot_bgcolor': '#ffffff', 'paper_bgcolor': '#ffffff',
                       'margin' : {'t':250, 'b':3, 'l':3, 'r':3}, "autosize": True
                      })

    return(fig)


#-- 3. Plot Wire Rope (Bulk Sumba)
def plot_wr1_s(value, wire_rope_c1):
    anchos=0.5
    color=[]
    position=[]
    for i in value:
        if i<=20 : color.append('#488f31')
        elif i<=40 : color.append('#88c580')
        elif i<=60 : color.append('#ffe600')
        elif i<=80 : color.append('#f7a258')
        else: color.append('#de425b')
    
    for i in value:
        if i<=50 : position.append('outside')
        else : position.append('inside')

    fig = go.Figure()
    fig.add_trace(go.Bar(x = wire_rope_c1['Crane 1'],
                         y = wire_rope_c1['Max'],
                         marker_color='#ededed',
                         hovertemplate='%{y:y}%',
                         width = anchos, name = 'Crane 1'))
    fig.add_trace(go.Bar(x = wire_rope_c1['Crane 1'], 
                         y = value,
                         text=value.apply(lambda x: '{0:1.0f}%'.format(x)),
                         textposition=position,
                         textfont_size=60,
                         textangle=0,
                         hovertemplate='%{y:y}%',
                         marker_color =color,
                         width = anchos, name = 'Crane 1'))
    fig.update_yaxes(visible=False)
    fig.update_layout(title = "<b>Crane 1</b>",
                      title_x=0.5,
                      barmode = 'overlay',
                      title_font_size = 78,
                      showlegend=False,
                      xaxis = dict(tickfont = dict(size=60)))
    fig.update_layout({'height':1200,'width':1700,'plot_bgcolor': '#ffffff', 'paper_bgcolor': '#ffffff',
                       'margin' : {'t':250, 'b':3, 'l':3, 'r':3}, "autosize": True
                      })

    return(fig)

############################################################################################################################

def plot_wr2_s(value, wire_rope_c2):
    anchos = 0.5
    color=[]
    position=[]
    for i in value:
        if i<=20 : color.append('#488f31')
        elif i<=40 : color.append('#88c580')
        elif i<=60 : color.append('#ffe600')
        elif i<=80 : color.append('#f7a258')
        else: color.append('#de425b')
    
    for i in value:
        if i<=50 : position.append('outside')
        else : position.append('inside')

    fig = go.Figure()
    fig.add_trace(go.Bar(x = wire_rope_c2['Crane 2'],
                         y = wire_rope_c2['Max'],
                         marker_color='#ededed',
                         hovertemplate='%{y:y}%',
                         width = anchos, name = 'Crane 2'))
    fig.add_trace(go.Bar(x = wire_rope_c2['Crane 2'], 
                         y = value,
                         text=value.apply(lambda x: '{0:1.0f}%'.format(x)),
                         textposition=position,
                         textfont_size=60,
                         textangle=0,
                         hovertemplate='%{y:y}%',
                         marker_color =color,
                         width = anchos, name = 'Crane 2'))
    fig.update_yaxes(visible=False)
    fig.update_layout(title = "<b>Crane 2</b>",
                      title_x=0.5,
                      barmode = 'overlay',
                      title_font_size = 78,
                      showlegend=False,
                      xaxis = dict(tickfont = dict(size=60)))
    fig.update_layout({'height':1200,'width':1700,'plot_bgcolor': '#ffffff', 'paper_bgcolor': '#ffffff',
                       'margin' : {'t':250, 'b':3, 'l':3, 'r':3}, "autosize": True
                      })

    return(fig)



#-- 3. Plot Wire Rope (Bulk Sumba)
def plot_wr1_n(value, wire_rope_c1):
    anchos=0.5
    color=[]
    position=[]
    for i in value:
        if i<=20 : color.append('#488f31')
        elif i<=40 : color.append('#88c580')
        elif i<=60 : color.append('#ffe600')
        elif i<=80 : color.append('#f7a258')
        else: color.append('#de425b')
    
    for i in value:
        if i<=50 : position.append('outside')
        else : position.append('inside')

    fig = go.Figure()
    fig.add_trace(go.Bar(x = wire_rope_c1['Crane 1'],
                         y = wire_rope_c1['Max'],
                         marker_color='#ededed',
                         hovertemplate='%{y:y}%',
                         width = anchos, name = 'Crane 1'))
    fig.add_trace(go.Bar(x = wire_rope_c1['Crane 1'], 
                         y = value,
                         text=value.apply(lambda x: '{0:1.0f}%'.format(x)),
                         textposition=position,
                         textfont_size=60,
                         textangle=0,
                         hovertemplate='%{y:y}%',
                         marker_color =color,
                         width = anchos, name = 'Crane 1'))
    fig.update_yaxes(visible=False)
    fig.update_layout(title = "<b>Crane 1</b>",
                      title_x=0.5,
                      barmode = 'overlay',
                      title_font_size = 78,
                      showlegend=False,
                      xaxis = dict(tickfont = dict(size=60)))
    fig.update_layout({'height':2500,'width':1700,'plot_bgcolor': '#ffffff', 'paper_bgcolor': '#ffffff',
                       'margin' : {'t':250, 'b':200, 'l':3, 'r':3}, "autosize": True
                      })

    return(fig)

#-- 4. Plot Conveyor Belt Replacement
def plot_cbr(value, bulk, conveyor_br):
    anchos = 0.5
    color=[]
    for i in value:
        if i<=20 : color.append('#488f31')
        elif i<=40 : color.append('#88c580')
        elif i<=60 : color.append('#ffe600')
        elif i<=80 : color.append('#f7a258')
        else: color.append('#de425b')

    fig = go.Figure()
    fig.add_trace(go.Bar(x = conveyor_br['Max'][::-1],
                         y = conveyor_br[bulk][::-1],
                         marker_color='#ededed',
                         width = anchos, name = 'Crane 1',
                         orientation='h'))
    fig.add_trace(go.Bar(x = value[::-1], 
                         y = conveyor_br[bulk][::-1],
                         text=value[::-1].apply(lambda x: '{0:1.0f}%'.format(x)),
                         textposition='outside',
                         textfont_size=40,
                         marker_color =color[::-1],
                         width = anchos, name = 'Crane 1',
                         orientation='h'))
    fig.update_xaxes(visible=False)
    fig.update_layout(title = "<b>Conveyor Belt Utilization</b> <br>FB 8K & CB 5K",
                      title_x=0.5,
                      barmode = 'overlay',
                      title_font_size = 80,
                      showlegend=False,
                      yaxis = dict(tickfont = dict(size=55)))
    fig.update_layout({'height':1000,'width':2440,'plot_bgcolor': '#ffffff', 'paper_bgcolor': '#ffffff',
                       'margin' : {'t':240,'b':3, 'l':3, 'r':3}, "autosize": True
                      })

    return(fig)


#-- 5. Plot Alternator OH
def plot_aoh(value, bulk, alternator_oh):
    anchos = 0.5
    color=[]
    position=[]
    for i in value:
        if i<=20 : color.append('#488f31')
        elif i<=40 : color.append('#88c580')
        elif i<=60 : color.append('#ffe600')
        elif i<=80 : color.append('#f7a258')
        else: color.append('#de425b')

    for i in value:
        if i<=5 : position.append('outside')
        else : position.append('inside')

    fig = go.Figure()
    fig.add_trace(go.Bar(x = alternator_oh['Max'][::-1],
                         y = alternator_oh[bulk][::-1],
                         marker_color='#ededed',
                         width = anchos, name = 'Crane 1',
                         orientation='h'))
    fig.add_trace(go.Bar(x = value[::-1], 
                         y = alternator_oh[bulk][::-1],
                         text=value[::-1].apply(lambda x: '{0:1.0f}%'.format(x)),
                         textposition=position[::-1],
                         textfont_size=55,
                         marker_color =color[::-1],
                         width = anchos, name = 'Crane 1',
                         orientation='h'))
    fig.update_xaxes(visible=False)
    fig.update_layout(title = "<b>Alternator Utilization</b> <br> To be cleaned and service every 5 years",
                      title_x=0.5,
                      barmode = 'overlay',
                      title_font_size = 80,
                      showlegend=False,
                      yaxis = dict(tickfont = dict(size=55)))
    fig.update_layout({'height':1000,'width':2000,'plot_bgcolor': '#ffffff', 'paper_bgcolor': '#ffffff',
                       'margin' : {'t':240,'b':3, 'l':3, 'r':3}, "autosize": True
                      })

    return(fig)

#-- 6. Plot Critical Spare
# plotly
def plot_cs(df):
    color=''
    if df['values'][0]<=(50) : color='#d9534f'
    elif df['values'][0]<=(80) : color='#f0ad4e'
    elif df['values'][0]<=100 : color='#02b875'

    fig = go.Figure(data=[go.Pie(labels=df['names'], values=df['values'], 
                                 hole=0.5,  
                                 textinfo='none',
                                 marker_colors=[color, '#ededed'])])
    fig.update_layout(title = "<b>Critical Spare</b>",
                      title_x=0.5,
                      title_font_size = 45,
                      showlegend=False,)
    fig.update_layout({'height':480,'width':400,'plot_bgcolor': '#ffffff', 'paper_bgcolor': '#ffffff',
                       'margin' : {'t':120, 'b':3, 'l':3, 'r':3}, "autosize": True})
    val=str(df['values'][0])
    fig.add_annotation(x=0.5, y=0.5, text=f'<b>{val}%</b>', 
                           font=dict(size=45), showarrow=False)
    return(fig)